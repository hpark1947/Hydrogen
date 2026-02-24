"""
수소 에너지 발표 PPT 생성 스크립트
- 18장 슬라이드, 비즈니스 프로페셔널 디자인
- python-pptx 기반
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
import os

# ── 상수 ──────────────────────────────────────────────
NAVY = RGBColor(0x1B, 0x3A, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_NAVY = RGBColor(0x2C, 0x5F, 0x8A)
TABLE_HEADER_BG = RGBColor(0x1B, 0x3A, 0x5C)
TABLE_ROW_ALT = RGBColor(0xE8, 0xF0, 0xF8)
TABLE_ROW_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_NAME = "맑은 고딕"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "수소에너지_발표자료.pptx")


# ── 유틸리티 함수 ─────────────────────────────────────
def set_font(run, size=18, bold=False, color=DARK_GRAY, name=FONT_NAME):
    """run에 폰트 스타일 적용"""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_background(slide, color=WHITE):
    """슬라이드 배경색 설정"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_navy_header_bar(slide):
    """상단 남색 바"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0),
        SLIDE_WIDTH, Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()


def add_green_accent_line(slide, top=Inches(1.2)):
    """녹색 포인트 라인"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), top,
        SLIDE_WIDTH, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()


def add_slide_number(slide, num, total=18):
    """슬라이드 번호"""
    txBox = slide.shapes.add_textbox(
        Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{num} / {total}"
    set_font(run, size=12, color=MEDIUM_GRAY)


def add_footer_line(slide):
    """하단 남색 가는 선"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(7.1),
        Inches(12.333), Inches(0.02)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()


def add_title_text(slide, title_text, left=Inches(0.6), top=Inches(0.2),
                   width=Inches(12), height=Inches(0.9)):
    """헤더 바 위에 제목 텍스트"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    set_font(run, size=34, bold=True, color=WHITE)
    return txBox


def add_body_textbox(slide, left=Inches(0.8), top=Inches(1.6),
                     width=Inches(11.7), height=Inches(5.2)):
    """본문 텍스트박스 생성"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_bullet_point(tf, text, level=0, size=16, bold=False, color=DARK_GRAY,
                     space_after=Pt(6), first=False):
    """불릿 포인트 추가"""
    if first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.level = level
    p.space_after = space_after
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    return p


def create_table(slide, rows, cols, left, top, width, height):
    """테이블 생성 후 shape 반환"""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    return table_shape, table


def set_cell(table, row, col, text, size=12, bold=False, color=DARK_GRAY,
             alignment=PP_ALIGN.LEFT, fill_color=None):
    """셀 텍스트 및 스타일 설정"""
    cell = table.cell(row, col)
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    # 셀 마진
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color


def style_header_row(table, headers, size=16):
    """테이블 헤더 행 스타일"""
    for i, h in enumerate(headers):
        set_cell(table, 0, i, h, size=size, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER, fill_color=TABLE_HEADER_BG)


def style_data_rows(table, data, start_row=1, size=14):
    """테이블 데이터 행 스타일"""
    for r_idx, row_data in enumerate(data):
        row_num = start_row + r_idx
        bg = TABLE_ROW_ALT if row_num % 2 == 0 else TABLE_ROW_WHITE
        for c_idx, val in enumerate(row_data):
            set_cell(table, row_num, c_idx, val, size=size, fill_color=bg)


def setup_content_slide(prs, title_text, slide_num):
    """표준 콘텐츠 슬라이드 설정 (배경 + 헤더바 + 제목 + 번호)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_background(slide, WHITE)
    add_navy_header_bar(slide)
    add_green_accent_line(slide)
    add_title_text(slide, title_text)
    add_footer_line(slide)
    add_slide_number(slide, slide_num)
    return slide


# ── 슬라이드 생성 함수 ───────────────────────────────

def slide_01_cover(prs):
    """슬라이드 1: 표지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, NAVY)

    # 상단 녹색 라인
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Inches(2.0),
        SLIDE_WIDTH, Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()

    # 메인 타이틀
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.4), Inches(10.3), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "수소 에너지의 현재와 미래"
    set_font(run, size=52, bold=True, color=WHITE)

    # 서브타이틀
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)
    run2 = p2.add_run()
    run2.text = "The Present and Future of Hydrogen Energy"
    set_font(run2, size=26, color=GREEN)

    # 하단 녹색 라인
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Inches(4.2),
        SLIDE_WIDTH, Inches(0.08)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = GREEN
    shape2.line.fill.background()

    # 발표자 정보
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(10.3), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p3 = tf2.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "수소 에너지 도서 3권 분석 기반 종합 발표"
    set_font(run3, size=22, color=RGBColor(0xBB, 0xBB, 0xBB))

    p4 = tf2.add_paragraph()
    p4.alignment = PP_ALIGN.CENTER
    p4.space_before = Pt(16)
    run4 = p4.add_run()
    run4.text = "2026.02"
    set_font(run4, size=18, color=RGBColor(0x99, 0x99, 0x99))


def slide_02_toc(prs):
    """슬라이드 2: 목차"""
    slide = setup_content_slide(prs, "목차  |  Table of Contents", 2)

    toc_items = [
        ("01", "왜 수소인가?"),
        ("02", "수소의 종류 (색깔별 분류)"),
        ("03", "수소 가치사슬"),
        ("04", "생산 기술 (수전해 비교)"),
        ("05", "저장 · 운송 기술"),
        ("06", "활용 분야 · 섹터 커플링"),
        ("07", "수소 vs 배터리"),
        ("08", "글로벌 수소 시장 전망"),
        ("09", "미국 수소 전략"),
        ("10", "유럽 수소 전략"),
        ("11", "중국 · 중동 수소 전략"),
        ("12", "한국 수소 정책"),
        ("13", "한국 기업 투자 현황"),
        ("14", "수소 경제 핵심 수치"),
        ("15", "도전과 과제"),
        ("16", "결론 및 시사점"),
    ]

    # 좌측 8개
    for i, (num, title) in enumerate(toc_items[:8]):
        top = Inches(1.6) + Inches(i * 0.62)
        txBox = slide.shapes.add_textbox(Inches(0.8), top, Inches(5.5), Inches(0.55))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run_num = p.add_run()
        run_num.text = f"  {num}   "
        set_font(run_num, size=20, bold=True, color=GREEN)
        run_title = p.add_run()
        run_title.text = title
        set_font(run_title, size=19, color=DARK_GRAY)

    # 우측 8개
    for i, (num, title) in enumerate(toc_items[8:]):
        top = Inches(1.6) + Inches(i * 0.62)
        txBox = slide.shapes.add_textbox(Inches(6.8), top, Inches(5.5), Inches(0.55))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run_num = p.add_run()
        run_num.text = f"  {num}   "
        set_font(run_num, size=20, bold=True, color=GREEN)
        run_title = p.add_run()
        run_title.text = title
        set_font(run_title, size=19, color=DARK_GRAY)


def slide_03_why_hydrogen(prs):
    """슬라이드 3: 왜 수소인가?"""
    slide = setup_content_slide(prs, "왜 수소인가?", 3)
    tf = add_body_textbox(slide)

    add_bullet_point(tf, "수소(H2)는 우주에서 가장 풍부한 원소이자 궁극의 청정 에너지원", first=True, size=22, bold=True, color=NAVY)
    add_bullet_point(tf, "연소 시 물(H2O)만 생성 — CO2 배출 제로", level=1, size=18)
    add_bullet_point(tf, "에너지 저장 매체이자 탄소 없는 연료 (이중 역할)", level=1, size=18)

    add_bullet_point(tf, "")
    add_bullet_point(tf, "재생에너지만으로는 탄소중립 불가능", size=22, bold=True, color=NAVY)
    add_bullet_point(tf, "태양광/풍력은 전력만 담당 — 총 에너지의 50%는 전기화 불가 (중공업, 장거리 수송, 계절저장)", level=1, size=18)
    add_bullet_point(tf, "배터리 한계: EV 배터리 ~450kg, 대형 트럭용은 ~4,500kg → 적재 불가", level=1, size=18)
    add_bullet_point(tf, "계절 저장(여름 잉여 → 겨울 수요): 배터리로는 물리적으로 불가능", level=1, size=18)

    add_bullet_point(tf, "")
    add_bullet_point(tf, "수소 + 전기 = '파워 커플'", size=22, bold=True, color=NAVY)
    add_bullet_point(tf, "전기: 가정, 경차량  |  수소: 중공업, 장거리 수송, 계절 저장, 화학 공정", level=1, size=18)
    add_bullet_point(tf, "섹터 커플링: 전기·수송·산업·난방을 하나로 연결하는 '만능 에너지 캐리어'", level=1, size=18)

    # 인용문 박스
    quote_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(6.0), Inches(10.3), Inches(0.8)
    )
    quote_box.fill.solid()
    quote_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF6, 0xEE)
    quote_box.line.fill.background()
    qtf = quote_box.text_frame
    qtf.word_wrap = True
    qp = qtf.paragraphs[0]
    qp.alignment = PP_ALIGN.CENTER
    qrun = qp.add_run()
    qrun.text = '"청정에너지는 더러운 에너지만큼 저렴해질 때까지 대세가 되지 않는다." — Marco Alvera'
    set_font(qrun, size=16, color=NAVY, bold=True)


def slide_04_hydrogen_types(prs):
    """슬라이드 4: 수소의 종류"""
    slide = setup_content_slide(prs, "수소의 종류  |  색깔별 분류", 4)

    headers = ["구분", "생산 방식", "CO2 배출", "현황/전망"]
    data = [
        ["그레이 (Grey)", "천연가스 수증기개질(SMR)", "~10kg CO2/kg H2", "현재 생산의 75%"],
        ["블루 (Blue)", "SMR + 탄소 포집·저장(CCS)", "대폭 감소", "현실적 전환 다리"],
        ["그린 (Green)", "재생에너지 수전해", "제로", "궁극적 목표, 비용 급감 중"],
        ["핑크 (Pink)", "원자력 수전해", "제로 (핵폐기물)", "원전 보유국 중심"],
        ["터콰이즈", "메탄 열분해", "고체 탄소 (CO2 無)", "연구 단계"],
    ]

    _, table = create_table(
        slide, len(data) + 1, len(headers),
        Inches(0.6), Inches(1.6), Inches(12.1), Inches(3.5)
    )
    # 열 너비
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(3.8)
    table.columns[2].width = Inches(2.8)
    table.columns[3].width = Inches(3.3)

    style_header_row(table, headers)
    style_data_rows(table, data)

    # 전환 로드맵
    tf = add_body_textbox(slide, top=Inches(5.4), height=Inches(1.5))
    add_bullet_point(tf, "전환 로드맵 (3권 공통 합의)", first=True, size=20, bold=True, color=NAVY)
    add_bullet_point(tf, "현재: 그레이 지배(75%)  →  2020~30s: 블루 (다리 역할)  →  2040~50: 그린 전면 확산", level=1, size=17)


def slide_05_value_chain(prs):
    """슬라이드 5: 수소 가치사슬"""
    slide = setup_content_slide(prs, "수소 가치사슬  |  Value Chain", 5)

    # 4단계 박스
    stages = [
        ("생산", "수전해(AWE/PEM/SOEC)\nSMR + CCUS\n메탄 열분해"),
        ("저장", "고압 압축 (700bar)\n액화 (-253°C)\n암모니아/LOHC\n금속수소화물"),
        ("운송", "파이프라인\n해상 탱커\n튜브 트레일러"),
        ("활용", "수송 (FCEV)\n발전 (연료전지)\n산업 (철강/석유화학)\n건물 (난방/냉방)"),
    ]

    colors = [NAVY, LIGHT_NAVY, RGBColor(0x22, 0x8B, 0x5B), GREEN]

    for i, (title, content) in enumerate(stages):
        left = Inches(0.5) + Inches(i * 3.15)
        # 메인 박스
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, Inches(1.8), Inches(2.8), Inches(4.0)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i]
        box.line.fill.background()

        # 제목
        title_box = slide.shapes.add_textbox(left + Inches(0.1), Inches(1.9), Inches(2.6), Inches(0.7))
        ttf = title_box.text_frame
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.alignment = PP_ALIGN.CENTER
        trun = tp.add_run()
        trun.text = title
        set_font(trun, size=26, bold=True, color=WHITE)

        # 내용
        content_box = slide.shapes.add_textbox(left + Inches(0.2), Inches(2.7), Inches(2.4), Inches(2.8))
        ctf = content_box.text_frame
        ctf.word_wrap = True
        for j, line in enumerate(content.split('\n')):
            cp = ctf.paragraphs[0] if j == 0 else ctf.add_paragraph()
            cp.space_after = Pt(6)
            crun = cp.add_run()
            crun.text = f"  {line}"
            set_font(crun, size=16, color=WHITE)

        # 화살표 (마지막 제외)
        if i < 3:
            arrow_left = left + Inches(2.85)
            arrow_box = slide.shapes.add_textbox(arrow_left, Inches(3.5), Inches(0.3), Inches(0.5))
            atf = arrow_box.text_frame
            ap = atf.paragraphs[0]
            ap.alignment = PP_ALIGN.CENTER
            arun = ap.add_run()
            arun.text = "→"
            set_font(arun, size=28, bold=True, color=NAVY)

    # 하단 설명
    tf = add_body_textbox(slide, top=Inches(6.0), height=Inches(0.8))
    add_bullet_point(tf, "2050 글로벌 수소 운송: 파이프라인 55% + 암모니아 해상운송 40%", first=True, size=16, color=MEDIUM_GRAY)
    add_bullet_point(tf, "수소 저장·운송 시장: $21.7B (2030) → $566B (2050)", size=16, color=MEDIUM_GRAY)


def slide_06_production(prs):
    """슬라이드 6: 생산 기술"""
    slide = setup_content_slide(prs, "수전해 생산 기술 비교", 6)

    headers = ["구분", "AWE (알칼라인)", "PEM (고분자전해질)", "SOEC (고체산화물)"]
    data = [
        ["성숙도", "가장 성숙 (상용)", "상용화", "실증/초기"],
        ["시장점유율", "60~75%", "22~25%", "~4% (성장 중)"],
        ["핵심 장점", "최저 비용, 귀금속 불요\n대규모에 최적", "빠른 응답, 유연 운전\n(0~160%), 컴팩트", "최고 효율 (800°C+)\n부식 없음"],
        ["핵심 단점", "느린 응답 속도", "귀금속 촉매 고가", "초기 단계, 고온 필요"],
        ["주요 성과", "글로벌 주력 기술", "2005년 이후 자본비 80% 하락", "네덜란드 2.6MW\nNASA 4MW 설치"],
    ]

    _, table = create_table(
        slide, len(data) + 1, len(headers),
        Inches(0.5), Inches(1.6), Inches(12.3), Inches(4.2)
    )
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(3.5)
    table.columns[2].width = Inches(3.5)
    table.columns[3].width = Inches(3.5)

    style_header_row(table, headers)
    style_data_rows(table, data)

    tf = add_body_textbox(slide, top=Inches(6.1), height=Inches(0.8))
    add_bullet_point(tf, "2024년 글로벌 수전해 설비용량: 2 GW  |  시스템 비용 목표: $250~500/kW", first=True, size=16, color=MEDIUM_GRAY)
    add_bullet_point(tf, "KAIST: 백금 무함유 PEM 수전해 기술, 단원자 귀금속 AEM 촉매 개발", size=16, color=MEDIUM_GRAY)


def slide_07_storage_transport(prs):
    """슬라이드 7: 저장·운송 기술"""
    slide = setup_content_slide(prs, "저장 · 운송 기술 비교", 7)

    headers = ["방식", "원리", "장점", "단점"]
    data = [
        ["고압 압축\n(350~700bar)", "고압 탱크 저장", "최저 비용, FCEV 적용", "폭발 위험, 중단거리"],
        ["액화 수소\n(-253°C)", "극저온 액화", "부피 1/800 축소\n71 kg/m³", "액화 에너지 소비 큼"],
        ["암모니아\n(NH3)", "수소+질소 합성\n-33°C 액화", "LH2 대비 1.7배 수소 밀도\n121 kg/m³, 기존 인프라", "분해 에너지 필요"],
        ["LOHC\n(MCH/톨루엔)", "유기 액체 캐리어", "상온 액체, 석유화학\n인프라 활용", "고온 탈수소화 필요"],
        ["금속수소화물", "금속에 수소 흡수", "가장 안전한 저장", "저장 용량 제한"],
    ]

    _, table = create_table(
        slide, len(data) + 1, len(headers),
        Inches(0.4), Inches(1.6), Inches(12.5), Inches(4.5)
    )
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(2.8)
    table.columns[2].width = Inches(4.0)
    table.columns[3].width = Inches(3.7)

    style_header_row(table, headers)
    style_data_rows(table, data)

    tf = add_body_textbox(slide, top=Inches(6.3), height=Inches(0.6))
    add_bullet_point(tf, "Snam(이탈리아): 기존 천연가스 배관의 70%가 수소 호환  |  배관 1km = 수소 12톤 저장 = 4만 가구 1일 전력", first=True, size=16, color=MEDIUM_GRAY)


def slide_08_applications(prs):
    """슬라이드 8: 활용 분야"""
    slide = setup_content_slide(prs, "활용 분야  |  섹터 커플링", 8)

    sectors = [
        ("수송", [
            "FCEV 승용차 (현대 넥쏘 등)",
            "수소 트럭 (XCIENT 등 장거리)",
            "수소 선박 (암모니아 연료)",
            "항공 (e-kerosene 합성연료)",
        ]),
        ("발전", [
            "연료전지 발전 (15GW 목표)",
            "가스터빈 수소혼소/전소",
            "분산형 전원 (건물/공장)",
            "계절 저장 → 전력 변환",
        ]),
        ("산업", [
            "수소환원제철 (CO2 95% 감축)",
            "석유화학 원료",
            "시멘트 (초고온 공정)",
            "반도체 공정용 수소",
        ]),
        ("건물", [
            "수소 보일러 (난방)",
            "가정용 연료전지 (에네팜)",
            "수소 냉방 시스템",
            "천연가스 배관 수소 혼입",
        ]),
    ]

    colors = [NAVY, LIGHT_NAVY, RGBColor(0x22, 0x8B, 0x5B), GREEN]

    for i, (sector_name, items) in enumerate(sectors):
        left = Inches(0.4) + Inches(i * 3.15)

        # 섹터 제목 박스
        title_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, Inches(1.7), Inches(2.9), Inches(0.7)
        )
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = colors[i]
        title_box.line.fill.background()
        ttf = title_box.text_frame
        tp = ttf.paragraphs[0]
        tp.alignment = PP_ALIGN.CENTER
        trun = tp.add_run()
        trun.text = sector_name
        set_font(trun, size=24, bold=True, color=WHITE)

        # 항목들
        content_box = slide.shapes.add_textbox(left + Inches(0.15), Inches(2.6), Inches(2.7), Inches(3.2))
        ctf = content_box.text_frame
        ctf.word_wrap = True
        for j, item in enumerate(items):
            cp = ctf.paragraphs[0] if j == 0 else ctf.add_paragraph()
            cp.space_after = Pt(8)
            crun = cp.add_run()
            crun.text = f"  {item}"
            set_font(crun, size=16, color=DARK_GRAY)

    # 하단 포인트
    tf = add_body_textbox(slide, top=Inches(6.0), height=Inches(0.8))
    add_bullet_point(tf, "핵심 개념: 섹터 커플링 — 수소가 전기·수송·산업·난방을 하나로 연결하는 '만능 에너지 캐리어'",
                     first=True, size=17, bold=True, color=NAVY)
    add_bullet_point(tf, "장거리 트럭(2040): 최대 단일 수소 수요 분야 (~80 Mtpa)  |  항공(2050): ~50 Mtpa (e-fuel 기반)",
                     size=16, color=MEDIUM_GRAY)


def slide_09_h2_vs_battery(prs):
    """슬라이드 9: 수소 vs 배터리"""
    slide = setup_content_slide(prs, "수소 vs 배터리  |  영역별 적합성", 9)

    headers = ["분야", "수소 유리", "배터리/전기 유리", "비고"]
    data = [
        ["승용차", "", "배터리 우세", "BEV 경제성 확립"],
        ["대형 트럭", "수소 우세", "", "H2: ~900kg vs 배터리: ~4,500kg"],
        ["선박", "수소(암모니아) 우세", "", "글로벌 배출 ~3%, 암모니아 연료"],
        ["항공", "수소(e-kerosene)", "", "Airbus ZEROe 프로그램"],
        ["철강", "수소환원제철 (H2-DRI)", "", "CO2 95% 감축, 차 가격 <1% 상승"],
        ["계절 저장", "유일한 해법", "", "여름 잉여 → 겨울 수요"],
        ["장거리 에너지 수송", "파이프라인 우세", "", "송전선보다 훨씬 저렴"],
        ["가정/단거리", "", "전기 우세", "그리드 전력이 효율적"],
    ]

    _, table = create_table(
        slide, len(data) + 1, len(headers),
        Inches(0.4), Inches(1.6), Inches(12.5), Inches(4.8)
    )
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(3.0)
    table.columns[2].width = Inches(3.0)
    table.columns[3].width = Inches(4.3)

    style_header_row(table, headers)
    style_data_rows(table, data)

    tf = add_body_textbox(slide, top=Inches(6.6), height=Inches(0.4))
    add_bullet_point(tf, "수소와 전기는 경쟁이 아닌 보완 — 전기화 불가 영역에서 수소가 핵심 역할",
                     first=True, size=18, bold=True, color=NAVY)


def slide_10_global_market(prs):
    """슬라이드 10: 글로벌 수소 시장"""
    slide = setup_content_slide(prs, "글로벌 수소 시장  |  규모와 전망", 10)

    # 시장 규모 테이블
    headers = ["지표", "수치"]
    data = [
        ["2024 글로벌 수소 수요", "~1억 톤 (99% 그레이)"],
        ["2024 그린수소 시장", "$79.8억"],
        ["2025 글로벌 수소경제", "$1,863억 (CAGR 4.49%)"],
        ["2030 그린수소 시장", "$605.6억 (CAGR 38.5%)"],
        ["2050 글로벌 수소경제", "$2.5~12조"],
        ["2050 인프라 투자 기회", "$15~20조"],
        ["2050 수소 수요 전망", "3.75~6.6억 톤 (현재 대비 5~6배)"],
        ["2050 청정수소 에너지 비중", "25% (BNEF)"],
    ]

    _, table = create_table(
        slide, len(data) + 1, len(headers),
        Inches(0.6), Inches(1.6), Inches(5.5), Inches(4.8)
    )
    table.columns[0].width = Inches(2.8)
    table.columns[1].width = Inches(2.7)

    style_header_row(table, headers)
    style_data_rows(table, data, size=14)

    # 비용 하락 추이 (우측)
    cost_title = slide.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(5.5), Inches(0.5))
    ctf = cost_title.text_frame
    cp = ctf.paragraphs[0]
    crun = cp.add_run()
    crun.text = "그린수소 비용 하락 추이 ($/kg)"
    set_font(crun, size=20, bold=True, color=NAVY)

    cost_headers = ["시기", "비용", "vs 화석연료"]
    cost_data = [
        ["~2010", "$24/kg", "15배"],
        ["~2020", "$4~5/kg", "~2배"],
        ["2024", "$3~8/kg", "1.5~6배"],
        ["2030 목표", "$1.5~3/kg", "경쟁력 확보"],
        ["2050 목표", "$0.7~1.5/kg", "화석연료보다 저렴"],
    ]

    _, cost_table = create_table(
        slide, len(cost_data) + 1, len(cost_headers),
        Inches(6.8), Inches(2.2), Inches(5.5), Inches(3.2)
    )
    cost_table.columns[0].width = Inches(1.5)
    cost_table.columns[1].width = Inches(1.8)
    cost_table.columns[2].width = Inches(2.2)

    style_header_row(cost_table, cost_headers)
    style_data_rows(cost_table, cost_data, size=14)

    # 하단 비용 목표
    tf = add_body_textbox(slide, top=Inches(5.7), height=Inches(1.2))
    add_bullet_point(tf, "주요 비용 목표", first=True, size=18, bold=True, color=NAVY)
    add_bullet_point(tf, "US DOE: $1/kg (2031)  |  EU: $1.5~3/kg (2030)  |  한국: 3,500원/kg (2030) → 2,500원/kg (2050)", level=1, size=16)
    add_bullet_point(tf, "Green Hydrogen Catapult (7개 글로벌 기업): $2/kg (2026), 25GW 수전해기, 500만톤/년, $1,100억 투자", level=1, size=16)


def slide_11_us_strategy(prs):
    """슬라이드 11: 미국 수소 전략"""
    slide = setup_content_slide(prs, "미국 수소 전략  |  IRA와 수소허브", 11)
    tf = add_body_textbox(slide)

    add_bullet_point(tf, "IRA (인플레이션 감축법) — Section 45V", first=True, size=22, bold=True, color=NAVY)
    add_bullet_point(tf, "청정수소 생산세액공제: CO2 배출량 기준 4단계, 최대 $3/kg", level=1, size=18)
    add_bullet_point(tf, "10년간 총 ~$130억 규모 지원", level=1, size=18)
    add_bullet_point(tf, "석유·가스 초강대국에서 수소 초강대국으로의 확장 전략", level=1, size=18)

    add_bullet_point(tf, "")
    add_bullet_point(tf, "DOE Hydrogen Shot", size=22, bold=True, color=NAVY)
    add_bullet_point(tf, "목표: $1/kg by 2031 (1-1-1 비전)", level=1, size=18)

    add_bullet_point(tf, "")
    add_bullet_point(tf, "수소허브 프로그램 (H2Hubs)", size=22, bold=True, color=NAVY)
    add_bullet_point(tf, "지역별 생산-저장-활용 통합 생태계 구축", level=1, size=18)
    add_bullet_point(tf, "생산지-소비지-인프라를 클러스터로 연결", level=1, size=18)

    # IRA 세액공제 테이블
    headers = ["CO2 배출 (kg CO2/kg H2)", "세액공제 ($/kg)"]
    data = [
        ["< 0.45", "$3.00"],
        ["0.45 ~ 1.5", "$1.00"],
        ["1.5 ~ 2.5", "$0.75"],
        ["2.5 ~ 4.0", "$0.60"],
    ]
    _, table = create_table(
        slide, len(data) + 1, len(headers),
        Inches(7.5), Inches(1.8), Inches(4.8), Inches(2.5)
    )
    table.columns[0].width = Inches(2.8)
    table.columns[1].width = Inches(2.0)
    style_header_row(table, headers, size=15)
    style_data_rows(table, data, size=14)


def slide_12_eu_strategy(prs):
    """슬라이드 12: 유럽 수소 전략"""
    slide = setup_content_slide(prs, "유럽 수소 전략  |  REPowerEU", 12)
    tf = add_body_textbox(slide)

    add_bullet_point(tf, "REPowerEU 목표 (2030)", first=True, size=22, bold=True, color=NAVY)
    add_bullet_point(tf, "그린수소 1,000만 톤 자체 생산 + 1,000만 톤 수입", level=1, size=18)
    add_bullet_point(tf, "수소 프로젝트 총 투자: $1,340억", level=1, size=18)
    add_bullet_point(tf, "러시아 가스 의존 탈피 + 산업 경쟁력 확보 (이중 동기)", level=1, size=18)

    add_bullet_point(tf, "")
    add_bullet_point(tf, "European Hydrogen Backbone (유럽 수소 배관망)", size=22, bold=True, color=NAVY)
    add_bullet_point(tf, "2030: 31,500km  →  2040: 57,600km", level=1, size=18)
    add_bullet_point(tf, "IPCEI Hy2Infra: 최대 69억 유로 공적 자금", level=1, size=18)
    add_bullet_point(tf, "독일 단독: 200억 유로 핵심 수소 네트워크 투자", level=1, size=18)

    add_bullet_point(tf, "")
    add_bullet_point(tf, "SoutH2 Corridor", size=22, bold=True, color=NAVY)
    add_bullet_point(tf, "북아프리카 → 이탈리아 → 오스트리아 → 독일, 3,300km 파이프라인", level=1, size=18)
    add_bullet_point(tf, "유럽 에너지 안보의 핵심 인프라로 부상", level=1, size=18)


def slide_13_china_mideast(prs):
    """슬라이드 13: 중국·중동 수소 전략"""
    slide = setup_content_slide(prs, "중국 · 중동 수소 전략", 13)

    # 중국 섹션
    tf = add_body_textbox(slide, width=Inches(5.5), height=Inches(5.0))
    add_bullet_point(tf, "중국 — 세계 최대 수소 생산국", first=True, size=21, bold=True, color=NAVY)
    add_bullet_point(tf, "생산: 3,500만 톤 (2023), 용량 4,900만 톤", level=1, size=17)
    add_bullet_point(tf, "수소 산업 규모 1조 위안 돌파", level=1, size=17)
    add_bullet_point(tf, "수전해기: 글로벌 설치용량 65%\n    확정 주문 75% 점유", level=1, size=17)
    add_bullet_point(tf, "신장 쿠차 프로젝트: 세계 최초\n    대규모 상업 그린수소 (1,000m³/hr)", level=1, size=17)
    add_bullet_point(tf, "그린수소 프로젝트: 600+건 진행 중", level=1, size=17)
    add_bullet_point(tf, "FCEV 목표: 2025년 5만 대", level=1, size=17)

    # 중동 섹션
    tf2 = add_body_textbox(slide, left=Inches(6.8), width=Inches(5.5), height=Inches(5.0))
    add_bullet_point(tf2, "중동 — 석유 수출국의 대전환", first=True, size=21, bold=True, color=NAVY)

    add_bullet_point(tf2, "")
    add_bullet_point(tf2, "사우디아라비아", size=18, bold=True, color=LIGHT_NAVY)
    add_bullet_point(tf2, "세계 최대 석유 수출국 → 수소 수출국 전환", level=1, size=17)
    add_bullet_point(tf2, "NEOM Helios: $50억, 세계 최대 그린수소", level=1, size=17)

    add_bullet_point(tf2, "")
    add_bullet_point(tf2, "UAE", size=18, bold=True, color=LIGHT_NAVY)
    add_bullet_point(tf2, "블루수소 + 암모니아 수출 확대", level=1, size=17)
    add_bullet_point(tf2, "아부다비: 연 20만톤 암모니아 → 한국 수출", level=1, size=17)

    add_bullet_point(tf2, "")
    add_bullet_point(tf2, "오만", size=18, bold=True, color=LIGHT_NAVY)
    add_bullet_point(tf2, "그린수소 허브 추진", level=1, size=17)
    add_bullet_point(tf2, "POSCO 컨소시엄: 47년 독점 개발권 확보", level=1, size=17)


def slide_14_korea_policy(prs):
    """슬라이드 14: 한국 수소 정책"""
    slide = setup_content_slide(prs, "한국 수소 정책  |  로드맵과 목표", 14)

    headers = ["정책/지표", "내용"]
    data = [
        ["2005년", "최초 '수소경제 마스터플랜' 수립"],
        ["2019 로드맵", "FCEV 620만대, 충전소 1,200개, 연료전지 15GW (2040)"],
        ["2021 기본계획", "2030: 390만톤 / 2050: 2,790만톤 청정수소"],
        ["세계 최초 수소법", "수소경제 전용 법률 제정 (글로벌 최초)"],
        ["청정수소 입찰제(CHPS)", "세계 최초 시행 (2024), 입찰가 477원/kWh"],
        ["2050 수소 비중", "최종에너지 수요의 21%"],
        ["수입 의존도", "81% (~2,200만톤) 수입 필요 (2050)"],
        ["국내 생산", "블루 200만톤 + 그린 300만톤 = 500만톤 (2050)"],
        ["경제적 효과", "연 43조원 부가가치, 42만개 일자리 (2040)"],
        ["충전소 현황", "385개 (2024) → 450개 (2025) → 660개 (2030)"],
        ["수소 배관", "410km 건설 중"],
    ]

    _, table = create_table(
        slide, len(data) + 1, len(headers),
        Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5)
    )
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(8.8)

    style_header_row(table, headers, size=16)
    style_data_rows(table, data, size=14)


def slide_15_korea_companies(prs):
    """슬라이드 15: 한국 기업 투자 현황"""
    slide = setup_content_slide(prs, "한국 기업 수소 투자 현황", 15)

    headers = ["기업/그룹", "2030 투자규모", "주요 수소 사업"]
    data = [
        ["SK그룹\n(SK E&S)", "~$120억\n(~12조원)", "보령 블루수소 25만톤/년(2026)\n인천 세계최대 수소액화플랜트 3만톤/년"],
        ["현대차그룹", "~$72억", "FCEV 리더 (넥쏘, XCIENT 수소트럭)\n수소 상용 모빌리티 밸류체인 통합"],
        ["포스코", "~$65억", "수소환원제철(HyREX)\n오만 그린수소 22만톤/년 (47년 독점)"],
        ["롯데케미칼", "6조원", "2030년까지 청정수소 120만톤 생산·공급"],
        ["두산퓨얼셀", "—", "정치용 연료전지 국내 M/S 1위\n미국 수출 확대, SOFC 라인업"],
        ["한화", "—", "해양 수소 연료전지 (200kW, DNV인증)\n수소 저장 사업"],
    ]

    _, table = create_table(
        slide, len(data) + 1, len(headers),
        Inches(0.4), Inches(1.5), Inches(12.5), Inches(4.8)
    )
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(8.3)

    style_header_row(table, headers, size=16)
    style_data_rows(table, data, size=14)

    tf = add_body_textbox(slide, top=Inches(6.5), height=Inches(0.5))
    add_bullet_point(tf, "5대 그룹 합계: 2030년까지 $380억 (약 43조원) 투자 계획",
                     first=True, size=18, bold=True, color=NAVY)


def slide_16_key_numbers(prs):
    """슬라이드 16: 수소 경제 핵심 수치"""
    slide = setup_content_slide(prs, "수소 경제 핵심 수치  |  Summary Stats", 16)

    headers = ["분류", "지표", "수치"]
    data = [
        ["시장", "2024 글로벌 수소 수요", "~1억 톤"],
        ["시장", "2050 수소 수요 전망", "3.75~6.6억 톤 (5~6배)"],
        ["시장", "2050 수소경제 규모", "$2.5~12조"],
        ["비용", "그린수소 비용 (2024)", "$3~8/kg"],
        ["비용", "그린수소 목표 (2050)", "$0.7~1.5/kg"],
        ["투자", "EU 수소 프로젝트", "$1,340억"],
        ["투자", "한국 5대 그룹", "$380억 (2030)"],
        ["인프라", "EU 수소 배관 (2040)", "57,600km"],
        ["인프라", "한국 충전소 (2024)", "385개"],
        ["한국", "FCEV 등록 (누적)", "19,270대 (글로벌 1위)"],
        ["한국", "정치용 연료전지", "1GW+ (세계 1/3 이상)"],
        ["한국", "2050 수소 에너지 비중", "최종 에너지의 21%"],
    ]

    _, table = create_table(
        slide, len(data) + 1, len(headers),
        Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5)
    )
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(4.5)
    table.columns[2].width = Inches(6.3)

    style_header_row(table, headers, size=16)
    style_data_rows(table, data, size=14)


def slide_17_challenges(prs):
    """슬라이드 17: 도전과 과제"""
    slide = setup_content_slide(prs, "도전과 과제", 17)

    # 4개 영역을 2x2 박스로
    challenges = [
        ("비용", NAVY, [
            "그린수소: 화석 대비 1.5~6배 비쌈",
            "에너지 변환 손실 60~70%",
            "청정수소 입찰가 477원/kWh\n    — 수익성 미확보",
        ]),
        ("인프라", LIGHT_NAVY, [
            "2050까지 $15조 투자 필요",
            "한국 충전소 385개 (2024)",
            "수소 액화플랜트\n    — 전량 외국 기술 의존",
        ]),
        ("실행 리스크", RGBColor(0x22, 0x8B, 0x5B), [
            "글로벌 프로젝트: 49→37 Mtpa 하향",
            "바인딩 계약 < 2 Mt/year",
            "전문가: 2030 목표의 10%만 달성 가능",
            "Stellantis, 수소차 프로그램 취소 (2025.7)",
        ]),
        ("정책·안전", GREEN, [
            "한국: 포지티브 규제 → 유연성 부족",
            "정권 교체 시 정책 일관성 리스크",
            "중국 수전해기 65~75% 지배\n    — 공급망 리스크",
            "수소 저장·충전 사고 60%+ 집중",
        ]),
    ]

    positions = [
        (Inches(0.4), Inches(1.6)),
        (Inches(6.6), Inches(1.6)),
        (Inches(0.4), Inches(4.2)),
        (Inches(6.6), Inches(4.2)),
    ]

    for (title, color, items), (left, top) in zip(challenges, positions):
        # 제목 바
        title_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, Inches(6.0), Inches(0.55)
        )
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = color
        title_shape.line.fill.background()
        ttf = title_shape.text_frame
        tp = ttf.paragraphs[0]
        tp.alignment = PP_ALIGN.CENTER
        trun = tp.add_run()
        trun.text = title
        set_font(trun, size=20, bold=True, color=WHITE)

        # 내용
        content_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.65), Inches(5.6), Inches(1.8))
        ctf = content_box.text_frame
        ctf.word_wrap = True
        for j, item in enumerate(items):
            cp = ctf.paragraphs[0] if j == 0 else ctf.add_paragraph()
            cp.space_after = Pt(4)
            crun = cp.add_run()
            crun.text = f"  {item}"
            set_font(crun, size=15, color=DARK_GRAY)

    # 인용문
    quote_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(6.8), Inches(10.3), Inches(0.5)
    )
    quote_box.fill.solid()
    quote_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF6, 0xEE)
    quote_box.line.fill.background()
    qtf = quote_box.text_frame
    qtf.word_wrap = True
    qp = qtf.paragraphs[0]
    qp.alignment = PP_ALIGN.CENTER
    qrun = qp.add_run()
    qrun.text = '"저탄소 수소는 실패하는 것이 아니라, 닷컴 시대 초기를 겪고 있는 것이다." — ERM Report (2024)'
    set_font(qrun, size=15, color=NAVY, bold=True)


def slide_18_conclusion(prs):
    """슬라이드 18: 결론 및 시사점"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, NAVY)

    # 상단 녹색 라인
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Inches(0.8),
        SLIDE_WIDTH, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.3), Inches(0.9))
    ttf = title_box.text_frame
    tp = ttf.paragraphs[0]
    tp.alignment = PP_ALIGN.CENTER
    trun = tp.add_run()
    trun.text = "결론 및 시사점"
    set_font(trun, size=42, bold=True, color=WHITE)

    # 핵심 메시지 3가지
    messages = [
        ("1", "수소는 탄소중립의 필수 퍼즐 조각",
         "전기화만으로는 순배출 제로 불가\n중공업·장거리 수송·계절 저장에 수소는 유일한 대안\n2050 수소 수요: 현재의 5~6배 (3.75~6.6억 톤)"),
        ("2", "그레이 → 블루 → 그린, 단계적 전환이 현실 경로",
         "그린수소 비용: $24/kg(2010) → $3~8(현재) → $1 이하(2050)\n블루수소가 2030년대까지 '다리' 역할 수행\n세계적 투자 물결: 수조 달러 규모 진행 중"),
        ("3", "한국: 기술 리더십 + 해외 공급망 + 제도적 시장 = 경쟁력",
         "FCEV·연료전지 세계 선도, 세계 최초 수소법·입찰제\n약점: 재생에너지 부족 → 81% 수입 의존 불가피\n핵심: 액화 기술 국산화 + 해외 파트너십 다변화"),
    ]

    for i, (num, title, desc) in enumerate(messages):
        top = Inches(2.2) + Inches(i * 1.6)

        # 번호 원
        num_box = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.0), top, Inches(0.6), Inches(0.6)
        )
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = GREEN
        num_box.line.fill.background()
        ntf = num_box.text_frame
        np_ = ntf.paragraphs[0]
        np_.alignment = PP_ALIGN.CENTER
        nrun = np_.add_run()
        nrun.text = num
        set_font(nrun, size=24, bold=True, color=WHITE)

        # 제목
        msg_title = slide.shapes.add_textbox(Inches(1.9), top - Inches(0.05), Inches(9.5), Inches(0.5))
        mtf = msg_title.text_frame
        mp = mtf.paragraphs[0]
        mrun = mp.add_run()
        mrun.text = title
        set_font(mrun, size=24, bold=True, color=GREEN)

        # 설명
        msg_desc = slide.shapes.add_textbox(Inches(1.9), top + Inches(0.5), Inches(9.5), Inches(1.0))
        dtf = msg_desc.text_frame
        dtf.word_wrap = True
        for j, line in enumerate(desc.split('\n')):
            dp = dtf.paragraphs[0] if j == 0 else dtf.add_paragraph()
            dp.space_after = Pt(3)
            drun = dp.add_run()
            drun.text = f"  {line}"
            set_font(drun, size=16, color=RGBColor(0xCC, 0xCC, 0xCC))

    # 하단 감사 메시지
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Inches(6.6),
        SLIDE_WIDTH, Inches(0.06)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = GREEN
    shape2.line.fill.background()

    thanks_box = slide.shapes.add_textbox(Inches(1.0), Inches(6.8), Inches(11.3), Inches(0.6))
    ttf2 = thanks_box.text_frame
    tp2 = ttf2.paragraphs[0]
    tp2.alignment = PP_ALIGN.CENTER
    trun2 = tp2.add_run()
    trun2.text = "감사합니다  |  Thank You"
    set_font(trun2, size=30, bold=True, color=WHITE)

    add_slide_number(slide, 18)


# ── 메인 실행 ─────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 18장 슬라이드 생성
    slide_01_cover(prs)
    slide_02_toc(prs)
    slide_03_why_hydrogen(prs)
    slide_04_hydrogen_types(prs)
    slide_05_value_chain(prs)
    slide_06_production(prs)
    slide_07_storage_transport(prs)
    slide_08_applications(prs)
    slide_09_h2_vs_battery(prs)
    slide_10_global_market(prs)
    slide_11_us_strategy(prs)
    slide_12_eu_strategy(prs)
    slide_13_china_mideast(prs)
    slide_14_korea_policy(prs)
    slide_15_korea_companies(prs)
    slide_16_key_numbers(prs)
    slide_17_challenges(prs)
    slide_18_conclusion(prs)

    prs.save(OUTPUT_PATH)
    print(f"PPT 생성 완료: {OUTPUT_PATH}")
    print(f"총 슬라이드 수: {len(prs.slides)}")


if __name__ == "__main__":
    main()
