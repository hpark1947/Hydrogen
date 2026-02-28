"""
수소자동차 시장 분석 PPT 생성 스크립트
- 28장 슬라이드, 비즈니스 프로페셔널 디자인
- python-pptx 기반
- 세계·한국·BEV 비교 종합분석
- 데이터 출처: 07_수소자동차_시장_현황과_전망.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 상수 ──────────────────────────────────────────────
NAVY = RGBColor(0x1B, 0x3A, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_NAVY = RGBColor(0x2C, 0x5F, 0x8A)
ACCENT_BLUE = RGBColor(0x34, 0x98, 0xDB)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_ORANGE = RGBColor(0xF3, 0x9C, 0x12)
TABLE_HEADER_BG = RGBColor(0x1B, 0x3A, 0x5C)
TABLE_ROW_ALT = RGBColor(0xE8, 0xF0, 0xF8)
TABLE_ROW_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_NAME = "맑은 고딕"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
TOTAL_SLIDES = 28

OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "수소자동차_시장분석_발표자료.pptx")


# ── 유틸리티 함수 ─────────────────────────────────────
def set_font(run, size=18, bold=False, color=DARK_GRAY, name=FONT_NAME):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_background(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_navy_header_bar(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_WIDTH, Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()


def add_green_accent_line(slide, top=Inches(1.2)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), top, SLIDE_WIDTH, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()


def add_slide_number(slide, num):
    txBox = slide.shapes.add_textbox(
        Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{num} / {TOTAL_SLIDES}"
    set_font(run, size=12, color=MEDIUM_GRAY)


def add_footer_line(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(7.1), Inches(12.333), Inches(0.02)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()


def add_title_text(slide, title_text, left=Inches(0.6), top=Inches(0.2),
                   width=Inches(12), height=Inches(0.9)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    set_font(run, size=34, bold=True, color=WHITE)
    return txBox


def add_body_textbox(slide, left=Inches(0.8), top=Inches(1.6),
                     width=Inches(11.7), height=Inches(5.2)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_bullet(tf, text, level=0, size=16, bold=False, color=DARK_GRAY,
               space_after=Pt(6), first=False):
    if first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.level = level
    p.space_after = space_after
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    return p


def set_cell(table, row, col, text, size=12, bold=False, color=DARK_GRAY,
             alignment=PP_ALIGN.LEFT, fill_color=None):
    cell = table.cell(row, col)
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color


def style_header_row(table, headers, size=14):
    for i, h in enumerate(headers):
        set_cell(table, 0, i, h, size=size, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER, fill_color=TABLE_HEADER_BG)


def style_data_rows(table, data, start_row=1, size=13):
    for r_idx, row_data in enumerate(data):
        row_num = start_row + r_idx
        bg = TABLE_ROW_ALT if row_num % 2 == 0 else TABLE_ROW_WHITE
        for c_idx, val in enumerate(row_data):
            set_cell(table, row_num, c_idx, val, size=size, fill_color=bg)


def setup_slide(prs, title_text, slide_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_navy_header_bar(slide)
    add_green_accent_line(slide)
    add_title_text(slide, title_text)
    add_footer_line(slide)
    add_slide_number(slide, slide_num)
    return slide


def add_colored_box(slide, left, top, width, height, fill_color, text,
                    text_size=13, text_color=WHITE, bold=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_font(run, size=text_size, bold=bold, color=text_color)
    return shape


# ── 슬라이드 생성 ─────────────────────────────────────

def slide_01_cover(prs):
    """표지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, NAVY)

    # 상단 녹색 라인
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Inches(2.0), SLIDE_WIDTH, Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()

    # 제목
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.4), Inches(10.3), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "수소자동차 시장 현황과 전망"
    set_font(run, size=46, bold=True, color=WHITE)

    # 부제
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.9), Inches(10.3), Inches(1.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "세계·한국·BEV 비교 종합분석"
    set_font(run2, size=24, bold=False, color=RGBColor(0xBB, 0xD5, 0xED))

    # 날짜
    txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.6))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "2026. 02"
    set_font(run3, size=18, color=RGBColor(0x88, 0xAA, 0xCC))

    # 하단 녹색 라인
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Inches(5.7), SLIDE_WIDTH, Inches(0.08)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = GREEN
    shape2.line.fill.background()


def slide_02_toc(prs):
    """목차"""
    slide = setup_slide(prs, "목차 (Table of Contents)", 2)
    sections = [
        ("Part 1", "세계 수소차 시장 현황", "글로벌 판매 추이, 국가별 보급, 주요 모델, 충전 인프라"),
        ("Part 2", "시장 전망 및 정책", "시장 전망, 정부 정책, 상용차·기타 모빌리티"),
        ("Part 3", "핵심 기술 발전", "스택·촉매 기술, 비용·저장 혁신"),
        ("Part 4", "한국 시장 심층분석", "시장 현황, 인프라, 정책, 기업 생태계, 글로벌 위상"),
        ("Part 5", "BEV vs FCEV 비교 및 결론", "사양·TCO·환경성, 시나리오 분석, 전략 권고"),
    ]
    colors = [NAVY, LIGHT_NAVY, ACCENT_BLUE, ACCENT_ORANGE, GREEN]

    for i, (part, title, desc) in enumerate(sections):
        y = Inches(1.6) + Inches(i * 1.05)
        add_colored_box(slide, Inches(0.8), y, Inches(1.5), Inches(0.8),
                        colors[i], part, text_size=16, bold=True)
        txBox = slide.shapes.add_textbox(Inches(2.6), y, Inches(9.5), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        set_font(run, size=20, bold=True, color=DARK_GRAY)
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = desc
        set_font(run2, size=14, color=MEDIUM_GRAY)


def slide_03_global_market(prs):
    """세계 수소차 시장 현황"""
    slide = setup_slide(prs, "세계 수소자동차 시장 현황", 3)

    # 테이블: 연도별 판매 추이
    tbl_shape = slide.shapes.add_table(7, 5, Inches(0.6), Inches(1.5),
                                       Inches(8.5), Inches(4.0))
    table = tbl_shape.table
    table.columns[0].width = Inches(1.2)
    table.columns[1].width = Inches(1.8)
    table.columns[2].width = Inches(1.5)
    table.columns[3].width = Inches(1.5)
    table.columns[4].width = Inches(2.5)
    style_header_row(table, ["연도", "글로벌 판매량", "전년 대비", "누적 등록", "주요 이슈"])
    data = [
        ["2020년", "~9,000대", "-", "~3만 대", "코로나19, 한국 수소법"],
        ["2021년", "~16,000대", "+78%", "~4.6만 대", "미라이 2세대 출시"],
        ["2022년", "~20,000대", "+25%", "~6.6만 대", "한국/중국 주도 성장"],
        ["2023년", "~14,000대", "-30%", "~8만 대", "보조금 변화, 경기 둔화"],
        ["2024년", "~16,011대", "+14%", "~9.6만 대", "현대차 회복, 미라이↓"],
        ["2025년(추정)", "~20,000대", "+25%", "~11.6만 대", "넥쏘 2세대 효과"],
    ]
    style_data_rows(table, data, size=12)

    # BEV 규모 격차 박스
    add_colored_box(slide, Inches(9.5), Inches(1.5), Inches(3.5), Inches(1.3),
                    ACCENT_RED, "BEV와의 규모 격차\n2024년 기준 약 850배\nFCEV ~0.02% vs BEV ~18%",
                    text_size=13)
    add_colored_box(slide, Inches(9.5), Inches(3.0), Inches(3.5), Inches(1.3),
                    NAVY, "FCEV 시장 특성\n연 1~2만 대 정체\n극소 니치(niche) 영역",
                    text_size=13)
    add_colored_box(slide, Inches(9.5), Inches(4.5), Inches(3.5), Inches(1.0),
                    GREEN, "BEV: 310만→1,380만 대\n5년간 4.5배 성장",
                    text_size=13)


def slide_04_country_status(prs):
    """국가별 보급 현황"""
    slide = setup_slide(prs, "주요 국가별 FCEV 보급 현황 (2025년)", 4)

    # 테이블
    tbl_shape = slide.shapes.add_table(7, 4, Inches(0.6), Inches(1.5),
                                       Inches(12.1), Inches(3.5))
    table = tbl_shape.table
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(1.6)
    table.columns[3].width = Inches(6.0)
    style_header_row(table, ["국가", "누적 등록 대수", "세계 비중", "특징"])
    data = [
        ["한국", "~40,000대", "~35%", "세계 유일 성장 시장, 넥쏘 단일 모델"],
        ["중국", "~30,000대", "~26%", "버스/트럭 상용차 중심, 부품 국산화율 70%+"],
        ["미국", "~12,283대", "~11%", "캘리포니아 편중, 충전소 51개, 정책 불확실"],
        ["일본", "~8,000대", "~7%", "미라이 중심, 판매 둔화(-39.1%)"],
        ["유럽", "~5,000대", "~4%", "상용차/버스 전환, BMW-도요타 협력"],
        ["기타", "~20,000대+", "~17%", "-"],
    ]
    style_data_rows(table, data, size=12)

    # 하단 불릿
    tf = add_body_textbox(slide, top=Inches(5.3), height=Inches(1.5))
    add_bullet(tf, "한국: 넥쏘 글로벌 판매 6,861대(+78.9%), 시장점유율 42.9% 세계 1위",
               first=True, size=14, bold=True, color=NAVY)
    add_bullet(tf, "중국: 에너지법에 수소 공식 분류, 연료전지 특허 글로벌 69% 장악, 비용 연 33% 하락",
               size=14, color=DARK_GRAY)
    add_bullet(tf, "미국: 트럼프 행정부 IRA 45V 단축(10년→2년), H2Hub 22억$ 취소, FY2026 수소 예산 $0 제안",
               size=14, color=ACCENT_RED)


def slide_05_key_models(prs):
    """넥쏘 2세대 vs 미라이"""
    slide = setup_slide(prs, "주요 모델: 현대 넥쏘 2세대 vs 도요타 미라이", 5)

    # 넥쏘 2세대 테이블
    tf_label1 = add_body_textbox(slide, left=Inches(0.6), top=Inches(1.35),
                                  width=Inches(5.5), height=Inches(0.3))
    add_bullet(tf_label1, "현대 넥쏘 2세대 (2025년 출시)", first=True, size=16, bold=True, color=NAVY)

    tbl1 = slide.shapes.add_table(7, 2, Inches(0.6), Inches(1.7),
                                   Inches(5.5), Inches(3.5))
    t1 = tbl1.table
    t1.columns[0].width = Inches(2.2)
    t1.columns[1].width = Inches(3.3)
    style_header_row(t1, ["항목", "사양"])
    data1 = [
        ["연료전지 시스템", "2.5세대 (110kW, 실사용 94kW)"],
        ["모터 출력", "201~255마력"],
        ["주행거리", "약 720km (기록 ~1,400km)"],
        ["충전 시간", "약 5분"],
        ["V2L", "외부 전원 공급 지원"],
        ["가격(한국)", "~7,200만 원 (보조금 후 ~3,950만 원)"],
    ]
    style_data_rows(t1, data1, size=12)

    # 미라이 2세대 테이블
    tf_label2 = add_body_textbox(slide, left=Inches(6.5), top=Inches(1.35),
                                  width=Inches(6.0), height=Inches(0.3))
    add_bullet(tf_label2, "도요타 미라이 2세대", first=True, size=16, bold=True, color=NAVY)

    tbl2 = slide.shapes.add_table(6, 2, Inches(6.5), Inches(1.7),
                                   Inches(6.0), Inches(3.0))
    t2 = tbl2.table
    t2.columns[0].width = Inches(2.2)
    t2.columns[1].width = Inches(3.8)
    style_header_row(t2, ["항목", "사양"])
    data2 = [
        ["연료전지 출력", "128kW"],
        ["주행거리", "850km+(WLTP), 기네스 1,360km"],
        ["구동 방식", "후륜구동(FR)"],
        ["가격", "$50,000~$67,000(미국)"],
        ["판매 추이", "2024년 -39.1% 감소"],
    ]
    style_data_rows(t2, data2, size=12)

    # 시장점유율 박스
    add_colored_box(slide, Inches(0.6), Inches(5.5), Inches(5.5), Inches(1.0),
                    GREEN, "넥쏘: 글로벌 점유율 42.9% (세계 1위)\n양산 — 유럽형 2025.8 / 북미형 2025.11",
                    text_size=14)
    add_colored_box(slide, Inches(6.5), Inches(5.5), Inches(6.0), Inches(1.0),
                    ACCENT_RED, "미라이: 점유율 7.3%로 하락\n도요타 3세대 스택(5.4 kW/L) 2026년 준비",
                    text_size=14)


def slide_06_other_models(prs):
    """기타 모델 & 중단 사례"""
    slide = setup_slide(prs, "기타 주요 모델 및 개발 중단/철수 사례", 6)

    # 기타 모델 테이블
    tf_label1 = add_body_textbox(slide, left=Inches(0.6), top=Inches(1.35),
                                  width=Inches(6.0), height=Inches(0.3))
    add_bullet(tf_label1, "기타 주요 FCEV 모델", first=True, size=15, bold=True, color=NAVY)

    tbl1 = slide.shapes.add_table(4, 3, Inches(0.6), Inches(1.7),
                                   Inches(6.0), Inches(2.0))
    t1 = tbl1.table
    t1.columns[0].width = Inches(2.0)
    t1.columns[1].width = Inches(1.3)
    t1.columns[2].width = Inches(2.7)
    style_header_row(t1, ["모델", "제조사", "특징"])
    data1 = [
        ["iX5 Hydrogen", "BMW", "도요타 FC, 295kW, 2028년 양산"],
        ["CR-V e:FCEV", "혼다", "GM 공동개발, EPA 435km"],
        ["XCIENT Fuel Cell", "현대", "대형트럭, 180kW 듀얼, 725km"],
    ]
    style_data_rows(t1, data1, size=12)

    # 중단 사례 테이블
    tf_label2 = add_body_textbox(slide, left=Inches(0.6), top=Inches(3.9),
                                  width=Inches(6.0), height=Inches(0.3))
    add_bullet(tf_label2, "개발 중단/철수 사례", first=True, size=15, bold=True, color=ACCENT_RED)

    tbl2 = slide.shapes.add_table(4, 3, Inches(0.6), Inches(4.2),
                                   Inches(6.0), Inches(2.0))
    t2 = tbl2.table
    t2.columns[0].width = Inches(1.8)
    t2.columns[1].width = Inches(2.0)
    t2.columns[2].width = Inches(2.2)
    style_header_row(t2, ["제조사", "상태", "이유"])
    data2 = [
        ["Stellantis", "2025년 중단", "인프라 부족, 비용"],
        ["Nikola", "Chapter 11 파산", "95대 리콜, 과대 약속"],
        ["Mercedes-Benz", "GLC F-CELL 단종", "승용 BEV 전략 집중"],
    ]
    style_data_rows(t2, data2, size=12)

    # 오른쪽 박스
    add_colored_box(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(1.5),
                    NAVY, "제조사 전략 양극화\nBEV 올인: 테슬라, BYD, 폭스바겐\nBEV+FCEV 병행: 현대, 도요타, BMW, 혼다",
                    text_size=14)
    add_colored_box(slide, Inches(7.0), Inches(3.5), Inches(5.8), Inches(1.5),
                    ACCENT_ORANGE, "Nikola 파산의 교훈\n수소 트럭 상용화의 어려움\n과대 약속의 위험성 시사",
                    text_size=14)
    add_colored_box(slide, Inches(7.0), Inches(5.3), Inches(5.8), Inches(1.0),
                    GREEN, "BMW-도요타 협력\n2028년 양산형 수소차 출시 계획",
                    text_size=14)


def slide_07_charging_infra(prs):
    """글로벌 충전 인프라"""
    slide = setup_slide(prs, "글로벌 수소 충전 인프라 현황", 7)

    # 충전소 현황 테이블
    tbl1 = slide.shapes.add_table(8, 3, Inches(0.6), Inches(1.5),
                                   Inches(6.5), Inches(3.8))
    t1 = tbl1.table
    t1.columns[0].width = Inches(2.0)
    t1.columns[1].width = Inches(1.5)
    t1.columns[2].width = Inches(3.0)
    style_header_row(t1, ["국가/지역", "충전소 수", "특징"])
    data1 = [
        ["중국", "~540개", "세계 최다, 상용차 중심"],
        ["한국", "~407기", "세계 2위, 2030년 660기 목표"],
        ["일본", "~160개", "2030년 1,000기 목표"],
        ["독일", "~100개+", "유럽 최대"],
        ["미국", "~51개", "캘리포니아 편중"],
        ["기타 유럽", "~120개+", "프랑스, 네덜란드 등"],
        ["전 세계 합계", "~1,200개", "EV 대비 1:2,000~3,000"],
    ]
    style_data_rows(t1, data1, size=12)

    # 경제성 비교 테이블
    tf_label = add_body_textbox(slide, left=Inches(7.5), top=Inches(1.35),
                                 width=Inches(5.5), height=Inches(0.3))
    add_bullet(tf_label, "충전소 경제성 비교", first=True, size=15, bold=True, color=NAVY)

    tbl2 = slide.shapes.add_table(4, 3, Inches(7.5), Inches(1.7),
                                   Inches(5.3), Inches(2.0))
    t2 = tbl2.table
    t2.columns[0].width = Inches(1.8)
    t2.columns[1].width = Inches(1.8)
    t2.columns[2].width = Inches(1.7)
    style_header_row(t2, ["항목", "수소 충전소", "EV 급속 충전소"])
    data2 = [
        ["설치 비용", "30~50억 원/기", "5,000만~1억 원/기"],
        ["비용 차이", "기준", "30~100배 저렴"],
        ["수익성", "대부분 만성 적자", "가동률 따라 흑자"],
    ]
    style_data_rows(t2, data2, size=12)

    # 핵심 메시지 박스
    add_colored_box(slide, Inches(7.5), Inches(4.0), Inches(5.3), Inches(1.2),
                    ACCENT_RED, "최대 장벽\n충전소 설치비 30~50억 원\nEV 대비 30~100배 비쌈",
                    text_size=14)
    add_colored_box(slide, Inches(7.5), Inches(5.4), Inches(5.3), Inches(1.0),
                    NAVY, "전 세계 ~1,200개\nEV 충전소(수백만 기) 대비 극히 미미",
                    text_size=14)


def slide_08_market_outlook(prs):
    """시장 전망"""
    slide = setup_slide(prs, "세계 수소차 시장 전망", 8)

    # 연료전지 시장 규모 테이블
    tf_label1 = add_body_textbox(slide, left=Inches(0.6), top=Inches(1.35),
                                  width=Inches(6.0), height=Inches(0.3))
    add_bullet(tf_label1, "연료전지 전체 시장 규모 전망", first=True, size=15, bold=True, color=NAVY)

    tbl1 = slide.shapes.add_table(5, 3, Inches(0.6), Inches(1.7),
                                   Inches(6.0), Inches(2.5))
    t1 = tbl1.table
    t1.columns[0].width = Inches(2.5)
    t1.columns[1].width = Inches(2.0)
    t1.columns[2].width = Inches(1.5)
    style_header_row(t1, ["조사기관", "2030년 전망", "CAGR"])
    data1 = [
        ["Grand View Research", "369.8억 달러", "27.1%"],
        ["MarketsandMarkets", "181.6억 달러", "26.3%"],
        ["Wissen Research", "210억 달러", "25.0%"],
        ["MarkNtel Advisors", "81.3억 달러", "20.4%"],
    ]
    style_data_rows(t1, data1, size=12)

    # FCEV 보급 전망 테이블
    tf_label2 = add_body_textbox(slide, left=Inches(0.6), top=Inches(4.35),
                                  width=Inches(6.0), height=Inches(0.3))
    add_bullet(tf_label2, "FCEV 보급 전망 — 기관별 시나리오", first=True, size=15, bold=True, color=NAVY)

    tbl2 = slide.shapes.add_table(6, 3, Inches(0.6), Inches(4.7),
                                   Inches(6.0), Inches(2.0))
    t2 = tbl2.table
    t2.columns[0].width = Inches(2.0)
    t2.columns[1].width = Inches(1.5)
    t2.columns[2].width = Inches(2.5)
    style_header_row(t2, ["기관", "시나리오", "2030년 FCEV 비중"])
    data2 = [
        ["IEA NZE", "넷제로", "전체 EV의 0.2~0.5%"],
        ["BloombergNEF", "기본", "신차의 0.22%"],
        ["McKinsey", "가속", "대형 상용차 15~25%"],
        ["IRENA", "1.5도", "2050년 장거리 30%"],
        ["Hydrogen Council", "가속", "500만 대+(2030)"],
    ]
    style_data_rows(t2, data2, size=11)

    # 오른쪽 박스
    add_colored_box(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(1.5),
                    GREEN, "상용차 — 진정한 성장 동력\nCAGR 47.7% (2025~2032)\n대형 트럭/버스가 핵심 영역",
                    text_size=14)
    add_colored_box(slide, Inches(7.0), Inches(3.5), Inches(5.8), Inches(1.5),
                    NAVY, "승용차 FCEV 점유율\n1% 미만에 머물 전망\n대형 상용차가 핵심 성장 영역",
                    text_size=14)
    add_colored_box(slide, Inches(7.0), Inches(5.3), Inches(5.8), Inches(1.3),
                    ACCENT_ORANGE, "McKinsey 전망\n2040년 장거리 트럭\n수소 소비 ~80 Mtpa\n모빌리티 최대 수요처",
                    text_size=13)


def slide_09_government_policy(prs):
    """주요국 정책 비교"""
    slide = setup_slide(prs, "주요국 정부 정책 및 보조금 비교", 9)

    # 정책 비교 테이블
    tbl_shape = slide.shapes.add_table(6, 4, Inches(0.6), Inches(1.5),
                                       Inches(12.1), Inches(3.0))
    table = tbl_shape.table
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(3.5)
    table.columns[2].width = Inches(3.0)
    table.columns[3].width = Inches(4.1)
    style_header_row(table, ["국가", "핵심 정책", "FCEV 목표", "보조금 수준"])
    data = [
        ["한국", "수소경제 로드맵, 수소법(세계 최초)", "2030년 18만 대", "~3,250만 원(세계 최고)"],
        ["일본", "수소기본전략(세계 최초, 2017)", "2030년 80만 대", "~200~300만 엔"],
        ["중국", "수소산업 중장기 계획", "2025년 5만 대", "시범도시 16억 위안/년"],
        ["미국", "IRA, H2Hubs", "정량 목표 부재", "45V: $3/kg 세액공제"],
        ["EU/독일", "REPowerEU, H2 Kernnetz", "-", "~7,000~9,000유로"],
    ]
    style_data_rows(table, data, size=12)

    # 정책 리스크 불릿
    tf = add_body_textbox(slide, top=Inches(4.8), height=Inches(2.0))
    add_bullet(tf, "정책 리스크 요인", first=True, size=16, bold=True, color=ACCENT_RED)
    add_bullet(tf, "미국: IRA 45V 기간 10년→2년 단축, H2Hub 2개 22억$ 취소, FY2026 수소 예산 $0 제안",
               level=1, size=13)
    add_bullet(tf, "한국: 청정수소발전 입찰시장 낙찰률 11.5% 참패, CHPS 폐기론 대두",
               level=1, size=13)
    add_bullet(tf, "EU: REPowerEU FID 전환율 4%, 수전해 목표 40GW 대비 운영 중 385MW(1%)",
               level=1, size=13)
    add_bullet(tf, "중국: 지방정부 보조금 변동, 과잉 투자 우려",
               level=1, size=13)


def slide_10_commercial_vehicles(prs):
    """상용차 (트럭/버스)"""
    slide = setup_slide(prs, "수소 상용차 시장 동향 — 트럭 & 버스", 10)

    # XCIENT 테이블
    tf_label1 = add_body_textbox(slide, left=Inches(0.6), top=Inches(1.35),
                                  width=Inches(6.5), height=Inches(0.3))
    add_bullet(tf_label1, "현대 XCIENT Fuel Cell (대형 트럭)", first=True, size=15, bold=True, color=NAVY)

    tbl1 = slide.shapes.add_table(7, 2, Inches(0.6), Inches(1.7),
                                   Inches(5.5), Inches(3.0))
    t1 = tbl1.table
    t1.columns[0].width = Inches(2.0)
    t1.columns[1].width = Inches(3.5)
    style_header_row(t1, ["항목", "사양"])
    data1 = [
        ["연료전지", "쌍 90kW 스택, 총 180kW"],
        ["주행거리", "1회 충전 725km"],
        ["수소 탱크", "10개, 총 68kg"],
        ["유럽 배치", "165대, 누적 2,000만km"],
        ["북미 배치", "63대, ~160만km"],
        ["수상", "TIME '2025 최고의 발명'"],
    ]
    style_data_rows(t1, data1, size=12)

    # 수소 버스 현황 박스
    add_colored_box(slide, Inches(6.5), Inches(1.5), Inches(6.3), Inches(1.0),
                    GREEN, "수소 버스 — 한국: 누적 2,066대(+277%), 2030년 21,200대 목표",
                    text_size=14)
    add_colored_box(slide, Inches(6.5), Inches(2.7), Inches(6.3), Inches(0.8),
                    LIGHT_NAVY, "중국: 세계 최대(수천 대) | 유럽: 2025 상반기 279대(+426%)",
                    text_size=14)

    # 상용차 성장 전망
    add_colored_box(slide, Inches(6.5), Inches(3.8), Inches(6.3), Inches(1.2),
                    NAVY, "상용차 시장 전망\nCAGR 47.7% (2025~2032)\nDeloitte: 2035년 대형 트럭 FCEV 15~25%",
                    text_size=14)

    # Nikola 교훈
    add_colored_box(slide, Inches(6.5), Inches(5.3), Inches(6.3), Inches(1.0),
                    ACCENT_RED, "교훈: Nikola 파산(2025.2, Chapter 11)\n수소 트럭 95대 리콜, 과대 약속의 위험",
                    text_size=13)


def slide_11_other_mobility(prs):
    """기타 모빌리티"""
    slide = setup_slide(prs, "기타 수소 모빌리티 — 선박·항공·열차·지게차·잠수함", 11)

    items = [
        ("선박", "Viking Libra 수소 크루즈선\n6MW PEM(2026년)\n한화에어로 200kW 선박용 FC", NAVY),
        ("항공", "ZeroAvia ZA600(600kW)\n2027년 인증 목표\nAirbus ZEROe(2035년)", LIGHT_NAVY),
        ("열차", "Alstom Coradia iLint\n세계 최초 수소열차\n2022년 독일 상업 운행", ACCENT_BLUE),
        ("지게차", "Amazon 15,000대+\nWalmart 9,500대+\n빠른 충전, 냉동창고 적합", GREEN),
        ("잠수함", "범한퓨얼셀 AIP용 PEMFC\n장보고-III급\n100% 국산화", ACCENT_ORANGE),
        ("트랙터/드론", "현대차 수소전기 트랙터\n수소 드론 스타트업\n배터리 대비 3~5배 체공", ACCENT_RED),
    ]

    for i, (title, desc, color) in enumerate(items):
        col = i % 3
        row = i // 3
        x = Inches(0.6 + col * 4.2)
        y = Inches(1.5 + row * 2.8)
        add_colored_box(slide, x, y, Inches(3.9), Inches(2.5),
                        color, f"■ {title}\n\n{desc}", text_size=13)


def slide_12_stack_technology(prs):
    """스택/촉매 기술"""
    slide = setup_slide(prs, "핵심 기술: 스택 출력밀도 & 촉매 혁신", 12)

    # 스택 출력밀도 테이블
    tf_label1 = add_body_textbox(slide, left=Inches(0.6), top=Inches(1.35),
                                  width=Inches(6.5), height=Inches(0.3))
    add_bullet(tf_label1, "스택 출력밀도 세대별 향상", first=True, size=15, bold=True, color=NAVY)

    tbl1 = slide.shapes.add_table(6, 4, Inches(0.6), Inches(1.7),
                                   Inches(6.5), Inches(2.5))
    t1 = tbl1.table
    t1.columns[0].width = Inches(1.8)
    t1.columns[1].width = Inches(2.0)
    t1.columns[2].width = Inches(1.5)
    t1.columns[3].width = Inches(1.2)
    style_header_row(t1, ["제조사", "세대", "출력밀도", "시기"])
    data1 = [
        ["현대차", "2.5세대(넥쏘 2)", "~3.5 kW/L", "2025년"],
        ["현대차", "3세대", "5.0+ kW/L", "2027년"],
        ["도요타", "2세대", "4.4 kW/L", "2020년"],
        ["도요타", "3세대", "5.4 kW/L", "2026년"],
        ["중국 SynStack", "GIII", "4.5+ kW/L", "2024년"],
    ]
    style_data_rows(t1, data1, size=12)

    # 현대차 3세대 스택 목표 테이블
    tf_label2 = add_body_textbox(slide, left=Inches(0.6), top=Inches(4.4),
                                  width=Inches(6.5), height=Inches(0.3))
    add_bullet(tf_label2, "현대차 3세대 스택 목표 (2027년)", first=True, size=15, bold=True, color=NAVY)

    tbl2 = slide.shapes.add_table(5, 4, Inches(0.6), Inches(4.7),
                                   Inches(6.5), Inches(2.0))
    t2 = tbl2.table
    t2.columns[0].width = Inches(1.5)
    t2.columns[1].width = Inches(1.7)
    t2.columns[2].width = Inches(1.7)
    t2.columns[3].width = Inches(1.6)
    style_header_row(t2, ["구분", "2세대(넥쏘1)", "2.5세대(넥쏘2)", "3세대(개발중)"])
    data2 = [
        ["시기", "2018년", "2025년", "2027년"],
        ["스택 출력", "95kW", "110kW", "100/200kW"],
        ["내구성", "~16만km", "-", "50만km+"],
        ["가격", "기준", "-", "50%+ 인하"],
    ]
    style_data_rows(t2, data2, size=12)

    # 촉매 기술 불릿
    tf = add_body_textbox(slide, left=Inches(7.5), top=Inches(1.5), width=Inches(5.3), height=Inches(5.0))
    add_bullet(tf, "촉매 기술 혁신", first=True, size=16, bold=True, color=NAVY)
    add_bullet(tf, "백금 사용량: 1990년대 대비 kW당 90% 감소", level=1, size=13)
    add_bullet(tf, "DOE 2025 목표: 0.10 g/kW 미만", level=1, size=13)
    add_bullet(tf, "PGM-free 촉매(Fe-N-C): 0.85 W/cm² 달성", level=1, size=13)
    add_bullet(tf, "KAIST: 백금-아연 나노입자로 1/3 절감(2025)", level=1, size=13)


def slide_13_cost_storage(prs):
    """비용/저장 기술"""
    slide = setup_slide(prs, "연료전지 비용 추이 & 수소 저장 기술", 13)

    # 비용 추이 테이블
    tf_label1 = add_body_textbox(slide, left=Inches(0.6), top=Inches(1.35),
                                  width=Inches(5.5), height=Inches(0.3))
    add_bullet(tf_label1, "연료전지 시스템 비용 추이", first=True, size=15, bold=True, color=NAVY)

    tbl1 = slide.shapes.add_table(4, 2, Inches(0.6), Inches(1.7),
                                   Inches(5.0), Inches(1.8))
    t1 = tbl1.table
    t1.columns[0].width = Inches(2.5)
    t1.columns[1].width = Inches(2.5)
    style_header_row(t1, ["시점", "비용"])
    data1 = [
        ["현재", "$80~100/kW"],
        ["2030년 DOE 목표", "$80/kW"],
        ["장기 목표", "$35/kW"],
    ]
    style_data_rows(t1, data1, size=13)

    # 저장 기술 불릿
    tf = add_body_textbox(slide, left=Inches(0.6), top=Inches(3.8), width=Inches(5.5), height=Inches(2.5))
    add_bullet(tf, "수소 저장 기술 현황", first=True, size=16, bold=True, color=NAVY)
    add_bullet(tf, "현재 표준: 700bar 고압 기체(Type IV 탱크)", level=1, size=14)
    add_bullet(tf, "효성첨단소재 TANSOM: 2028년까지 1조 원 투자", level=1, size=14)
    add_bullet(tf, "일진하이솔루스/코오롱스페이스웍스: Type IV 수소 탱크", level=1, size=14)

    # 오른쪽 핵심 박스
    add_colored_box(slide, Inches(6.5), Inches(1.5), Inches(6.3), Inches(1.5),
                    GREEN, "비용 혁신 로드맵\n$80~100/kW → $35/kW\n가격이 내연기관 수준에 도달하면\nFCEV 시장 급성장 가능",
                    text_size=14)
    add_colored_box(slide, Inches(6.5), Inches(3.3), Inches(6.3), Inches(1.3),
                    NAVY, "중국 비용 혁신\n비용 연간 33% 하락\n부품 국산화율 70%+\n급격한 가격 경쟁력 확보",
                    text_size=14)
    add_colored_box(slide, Inches(6.5), Inches(4.9), Inches(6.3), Inches(1.3),
                    ACCENT_ORANGE, "탄소섬유 국산화\n효성 TANSOM — 고압 수소탱크 핵심 소재\n2028년까지 1조 원 투자 계획",
                    text_size=14)


def slide_14_korea_market(prs):
    """한국 시장 현황"""
    slide = setup_slide(prs, "한국 수소자동차 시장 현황", 14)

    # 연도별 판매 테이블
    tbl_shape = slide.shapes.add_table(8, 5, Inches(0.6), Inches(1.5),
                                       Inches(8.0), Inches(3.8))
    table = tbl_shape.table
    table.columns[0].width = Inches(1.2)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(1.3)
    table.columns[3].width = Inches(1.5)
    table.columns[4].width = Inches(2.5)
    style_header_row(table, ["연도", "신규 판매", "전년 대비", "누적 등록", "비고"])
    data = [
        ["2019년", "~4,200대", "-", "~5,000대", "넥쏘 본격 판매 시작"],
        ["2020년", "~5,800대", "+38%", "~10,000대", "코로나에도 성장"],
        ["2021년", "~8,500대", "+47%", "~19,000대", "최대 판매 기록"],
        ["2022년", "~10,800대", "+27%", "~29,000대", "보조금 축소에도 성장"],
        ["2023년", "~3,800대", "-65%", "~33,000대", "넥쏘 1세대 노후화"],
        ["2024년", "~3,700대", "-3%", "~38,000대", "2세대 출시 대기"],
        ["2025년", "~6,861대", "+78.9%", "~40,000대", "넥쏘 2세대 효과"],
    ]
    style_data_rows(table, data, size=11)

    # 오른쪽 박스
    add_colored_box(slide, Inches(9.0), Inches(1.5), Inches(3.8), Inches(1.3),
                    ACCENT_RED, "로드맵 목표 달성률\n2022년 목표 8.1만 대 대비\n실적 ~2.9만 대\n달성률 약 23%",
                    text_size=14)
    add_colored_box(slide, Inches(9.0), Inches(3.0), Inches(3.8), Inches(1.3),
                    GREEN, "넥쏘 글로벌 판매\n6,861대(+78.9%)\n한국 6,802대 / 해외 ~59대\n해외 비중 극히 제한적",
                    text_size=14)
    add_colored_box(slide, Inches(9.0), Inches(4.5), Inches(3.8), Inches(1.3),
                    NAVY, "해외 판매 부진 이유\n충전 인프라 절대 부족\nBEV 대비 가격 경쟁력 열위\n캘리포니아 외 보급 전무",
                    text_size=13)


def slide_15_korea_infra_crisis(prs):
    """한국 충전 인프라"""
    slide = setup_slide(prs, "한국 수소 충전 인프라 — 현황 및 수익성 위기", 15)

    # 충전소 현황 테이블
    tf_label1 = add_body_textbox(slide, left=Inches(0.6), top=Inches(1.35),
                                  width=Inches(5.5), height=Inches(0.3))
    add_bullet(tf_label1, "충전소 현황 및 목표", first=True, size=15, bold=True, color=NAVY)

    tbl1 = slide.shapes.add_table(6, 3, Inches(0.6), Inches(1.7),
                                   Inches(5.5), Inches(2.5))
    t1 = tbl1.table
    t1.columns[0].width = Inches(1.5)
    t1.columns[1].width = Inches(1.5)
    t1.columns[2].width = Inches(2.5)
    style_header_row(t1, ["시점", "충전소 수", "비고"])
    data1 = [
        ["2019년", "~30기", "초기 단계"],
        ["2022년", "~170기", "목표 310기 대비 55%"],
        ["2024년 말", "386기", "목표 385기 달성"],
        ["2025년 3월", "407기", "운영 중"],
        ["2030년 목표", "660기", "-"],
    ]
    style_data_rows(t1, data1, size=12)

    # 수익성 위기 테이블
    tf_label2 = add_body_textbox(slide, left=Inches(0.6), top=Inches(4.4),
                                  width=Inches(5.5), height=Inches(0.3))
    add_bullet(tf_label2, "충전소 수익성 위기", first=True, size=15, bold=True, color=ACCENT_RED)

    tbl2 = slide.shapes.add_table(7, 2, Inches(0.6), Inches(4.7),
                                   Inches(5.5), Inches(2.0))
    t2 = tbl2.table
    t2.columns[0].width = Inches(2.5)
    t2.columns[1].width = Inches(3.0)
    style_header_row(t2, ["항목", "수치"])
    data2 = [
        ["1기 설치비", "30~50억 원"],
        ["평균 가동률", "20~25%"],
        ["일 평균 충전", "4대에 불과"],
        ["흑자 충전소", "전국 7곳뿐"],
        ["수소 판매가(2025.7)", "10,239원/kg"],
        ["하이넷 4년 누적 적자", "166억 원"],
    ]
    style_data_rows(t2, data2, size=12)

    # 충전 인프라 기업 박스
    add_colored_box(slide, Inches(6.5), Inches(1.5), Inches(6.3), Inches(1.5),
                    NAVY, "충전 인프라 주요 기업\n효성중공업: 시장점유율 1위, 액화수소충전소 21개 계획\nSK E&S: 2026년까지 30개소(현재 17개소)\n하이넷: 만성 적자(4년 166억 원)",
                    text_size=13)
    add_colored_box(slide, Inches(6.5), Inches(3.3), Inches(6.3), Inches(1.5),
                    ACCENT_RED, "수소 가격 지속 인상\n2021년 8,000원/kg → 2025년 10,239원/kg\nkm당 비용: FCEV 100~120원 vs BEV 40~60원\nFCEV가 BEV의 약 2배",
                    text_size=13)
    add_colored_box(slide, Inches(6.5), Inches(5.1), Inches(6.3), Inches(1.3),
                    ACCENT_ORANGE, "BEV 충전기 수만 기 vs 수소 407기\n비율 1:100 이상\n인프라 격차 해소가 최대 과제",
                    text_size=14)


def slide_16_korea_policy(prs):
    """한국 정책/보조금"""
    slide = setup_slide(prs, "한국 수소 정책 프레임워크 & 보조금 비교", 16)

    # 정책 프레임워크 테이블
    tf_label1 = add_body_textbox(slide, left=Inches(0.6), top=Inches(1.35),
                                  width=Inches(6.5), height=Inches(0.3))
    add_bullet(tf_label1, "정부 정책 프레임워크", first=True, size=15, bold=True, color=NAVY)

    tbl1 = slide.shapes.add_table(6, 3, Inches(0.6), Inches(1.7),
                                   Inches(6.5), Inches(2.5))
    t1 = tbl1.table
    t1.columns[0].width = Inches(2.5)
    t1.columns[1].width = Inches(1.2)
    t1.columns[2].width = Inches(2.8)
    style_header_row(t1, ["정책", "시기", "핵심 내용"])
    data1 = [
        ["수소경제 활성화 로드맵", "2019.1", "수소차/연료전지 중심 초기 전략"],
        ["수소법(세계 최초)", "2020.2", "수소경제 육성 및 안전관리"],
        ["제1차 수소경제 이행 계획", "2021.11", "청정수소 중심 전환, 탄소중립"],
        ["청정수소발전 입찰시장", "2024.5", "세계 최초(낙찰률 11.5% 실패)"],
        ["CHPS 제도", "2024", "수소발전 의무화(폐기론 대두)"],
    ]
    style_data_rows(t1, data1, size=11)

    # 보조금 비교 테이블
    tf_label2 = add_body_textbox(slide, left=Inches(0.6), top=Inches(4.4),
                                  width=Inches(6.5), height=Inches(0.3))
    add_bullet(tf_label2, "보조금 비교: 넥쏘 vs BEV", first=True, size=15, bold=True, color=NAVY)

    tbl2 = slide.shapes.add_table(6, 3, Inches(0.6), Inches(4.7),
                                   Inches(6.5), Inches(2.0))
    t2 = tbl2.table
    t2.columns[0].width = Inches(2.0)
    t2.columns[1].width = Inches(2.3)
    t2.columns[2].width = Inches(2.2)
    style_header_row(t2, ["항목", "넥쏘 2세대(FCEV)", "EV6 등(BEV)"])
    data2 = [
        ["차량 가격", "~7,200만 원", "~5,000만 원"],
        ["국비 보조금", "~2,250만 원", "~780만 원"],
        ["지자체 보조금", "~1,000만 원", "~400만 원"],
        ["실구매가", "~3,950만 원", "~3,820만 원"],
        ["보조금 없을 때", "7,200만 원", "5,000만 원"],
    ]
    style_data_rows(t2, data2, size=12)

    # 오른쪽 핵심 박스
    add_colored_box(slide, Inches(7.5), Inches(1.5), Inches(5.3), Inches(1.5),
                    NAVY, "2030/2050 수치 목표\n수소차 18만 대 / 충전소 660기(2030)\n수소차 620만 대 / 충전소 1,200기(2040)\n수소 공급 2,790만 톤(2050)",
                    text_size=13)
    add_colored_box(slide, Inches(7.5), Inches(3.3), Inches(5.3), Inches(1.5),
                    ACCENT_RED, "보조금 의존 구조\n보조금 없으면 수소차는\nBEV 대비 44% 비싸\n시장 자생력 부재",
                    text_size=15)
    add_colored_box(slide, Inches(7.5), Inches(5.1), Inches(5.3), Inches(1.3),
                    ACCENT_ORANGE, "R&D 예산 감소\n2023년 3,339억 원\n→ 2025년 2,611억 원\n2년 연속 감소",
                    text_size=14)


def slide_17_korea_investment(prs):
    """대기업 투자"""
    slide = setup_slide(prs, "한국 대기업 그룹 수소 투자 규모 (2030년까지)", 17)

    # 투자 규모 테이블
    tbl_shape = slide.shapes.add_table(7, 3, Inches(0.6), Inches(1.5),
                                       Inches(8.0), Inches(3.8))
    table = tbl_shape.table
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(4.0)
    style_header_row(table, ["기업 그룹", "투자 규모", "주요 분야"])
    data = [
        ["SK 그룹", "18.5조 원", "수소 생산/유통/충전/SOFC 발전"],
        ["현대차그룹", "11.1조 원", "FCEV, HTWO 시스템, 인프라"],
        ["포스코 그룹", "~10조 원", "수소환원제철, 수소 생산"],
        ["한화 그룹", "수조 원", "수전해, 그린수소, 태양광"],
        ["효성 그룹", "1조 원+", "액화수소, 탄소섬유, 충전시스템"],
        ["5개 그룹 합산", "~43.4조 원", "수소 밸류체인 전 분야"],
    ]
    style_data_rows(table, data, size=13)

    # HTWO 브랜드 박스
    add_colored_box(slide, Inches(9.0), Inches(1.5), Inches(3.8), Inches(2.0),
                    NAVY, "현대 HTWO 전략\n2030년 70만 기 FC 판매 목표\n중국 광저우 해외 첫 공장\n2025.6 현대모비스 FC사업 인수\nR&D→생산 원스톱 체제",
                    text_size=13)
    add_colored_box(slide, Inches(9.0), Inches(3.8), Inches(3.8), Inches(1.5),
                    GREEN, "HTWO Energy Savannah\nClass-8 대형트럭 전용\n수소+전기 복합 충전 스테이션\n(업계 최초)",
                    text_size=13)
    add_colored_box(slide, Inches(9.0), Inches(5.5), Inches(3.8), Inches(1.0),
                    ACCENT_ORANGE, "5대 그룹 총 ~43.4조 원\n수소 밸류체인 전 분야 투자",
                    text_size=14)


def slide_18_korea_supply_chain(prs):
    """밸류체인/국산화"""
    slide = setup_slide(prs, "한국 수소차 밸류체인 & 핵심부품 국산화 현황", 18)

    # 밸류체인 불릿
    tf = add_body_textbox(slide, left=Inches(0.6), top=Inches(1.5), width=Inches(6.0), height=Inches(3.0))
    add_bullet(tf, "밸류체인별 주요 기업", first=True, size=16, bold=True, color=NAVY)
    add_bullet(tf, "수소 생산: SK E&S(블루수소, 보령 연 25만톤), 한화솔루션(수전해)",
               level=1, size=12)
    add_bullet(tf, "수소 저장: 효성첨단소재(TANSOM), 일진하이솔루스(Type IV), 가드넥",
               level=1, size=12)
    add_bullet(tf, "FC 스택: 현대차(PEMFC 차량용), 두산(PAFC/SOFC), 블룸SK(SOFC)",
               level=1, size=12)
    add_bullet(tf, "핵심 부품: 코오롱(전해질막, 국내 유일), 케이퓨얼셀(분리판, 원가 70%↓)",
               level=1, size=12)
    add_bullet(tf, "핵심 부품: 에프씨엠티(MEA 연 20만장), 제이앤티지(GDL 최초 상용화)",
               level=1, size=12)
    add_bullet(tf, "스타트업: 에스퓨얼셀(건물용 50%+), 미코파워(독자 SOFC), 엘켐텍",
               level=1, size=12)

    # 국산화 현황 테이블
    tf_label = add_body_textbox(slide, left=Inches(0.6), top=Inches(4.7),
                                 width=Inches(6.0), height=Inches(0.3))
    add_bullet(tf_label, "핵심부품 국산화율", first=True, size=15, bold=True, color=ACCENT_RED)

    tbl = slide.shapes.add_table(7, 3, Inches(0.6), Inches(5.0),
                                  Inches(6.0), Inches(1.8))
    t = tbl.table
    t.columns[0].width = Inches(2.0)
    t.columns[1].width = Inches(1.5)
    t.columns[2].width = Inches(2.5)
    style_header_row(t, ["부품", "국산화율", "평가"])
    data = [
        ["전해질막(PEM)", "10~20%", "심각 미흡"],
        ["촉매(Pt/C)", "10~15%", "심각 미흡"],
        ["MEA", "20~30%", "미흡"],
        ["GDL", "50~60%", "양호"],
        ["분리판", "60~70%", "양호"],
        ["BOP 부품", "70~80%", "양호"],
    ]
    style_data_rows(t, data, size=11)

    # 오른쪽 핵심 박스
    add_colored_box(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(1.5),
                    ACCENT_RED, "핵심 소재 격차\n전해질막(10~20%), 촉매(10~15%)\n선도국 대비 5~7년 격차\n공급망 리스크 노출",
                    text_size=14)
    add_colored_box(slide, Inches(7.0), Inches(3.3), Inches(5.8), Inches(1.5),
                    NAVY, "중국의 급추격\n부품 국산화율 70%+\n연료전지 특허 글로벌 69% 장악\n비용 연 33% 하락",
                    text_size=14)
    add_colored_box(slide, Inches(7.0), Inches(5.1), Inches(5.8), Inches(1.3),
                    GREEN, "범한퓨얼셀\n잠수함 AIP용 PEMFC\n100% 국산화 성공",
                    text_size=14)


def slide_19_korea_global_status(prs):
    """글로벌 위상"""
    slide = setup_slide(prs, "한국의 글로벌 위상 — 분야별 순위", 19)

    # 순위 테이블
    tbl_shape = slide.shapes.add_table(8, 3, Inches(0.6), Inches(1.5),
                                       Inches(8.0), Inches(3.8))
    table = tbl_shape.table
    table.columns[0].width = Inches(3.0)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(3.0)
    style_header_row(table, ["분야", "한국의 위상", "비고"])
    data = [
        ["수소전기차 판매", "세계 1위(현대차)", "점유율 42.9%"],
        ["수소차 보유 대수", "세계 2위", "누적 ~4만 대"],
        ["수소 대형트럭 해외 배치", "세계 1위", "XCIENT 228대+"],
        ["수소버스 보급", "세계 최상위", "2,066대"],
        ["발전용 연료전지 설치량", "세계 1위", "1,036MW(1GW 돌파)"],
        ["수소법 제정", "세계 최초", "2020년"],
        ["수소충전소 수", "세계 2위", "407기"],
    ]
    style_data_rows(table, data, size=13)

    # 구조적 과제 박스
    add_colored_box(slide, Inches(9.0), Inches(1.5), Inches(3.8), Inches(1.3),
                    ACCENT_RED, "충전 인프라 부족\n407기 vs BEV 수만 기\nkm당 비용 BEV의 2배",
                    text_size=13)
    add_colored_box(slide, Inches(9.0), Inches(3.0), Inches(3.8), Inches(1.3),
                    ACCENT_ORANGE, "그린수소 확보 문제\n90%+ 그레이수소\n2030년 청정수소 74% 수입",
                    text_size=13)
    add_colored_box(slide, Inches(9.0), Inches(4.5), Inches(3.8), Inches(1.0),
                    NAVY, "경제성 미확보\n보조금 없이 자생력 부재\nR&D 예산 2년 연속 감소",
                    text_size=13)
    add_colored_box(slide, Inches(9.0), Inches(5.7), Inches(3.8), Inches(1.0),
                    LIGHT_NAVY, "BEV 경쟁 압박\n승용: BEV 20만 vs FCEV 6,800\n비율 약 1:30",
                    text_size=13)


def slide_20_bev_vs_fcev_specs(prs):
    """사양 비교"""
    slide = setup_slide(prs, "BEV vs FCEV — 핵심 사양 비교", 20)

    # 비교 테이블
    tbl_shape = slide.shapes.add_table(9, 4, Inches(0.6), Inches(1.5),
                                       Inches(8.5), Inches(4.5))
    table = tbl_shape.table
    table.columns[0].width = Inches(2.8)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(2.0)
    table.columns[3].width = Inches(1.7)
    style_header_row(table, ["비교 항목", "FCEV", "BEV", "유리한 쪽"])
    data = [
        ["에너지 효율(WtW)", "25~35%", "70~90%", "BEV"],
        ["1회 충전 주행거리", "500~720km+", "300~600km", "FCEV"],
        ["충전 시간", "3~5분", "급속 20~40분", "FCEV"],
        ["차량 중량", "상대적 경량", "배터리 400~800kg", "FCEV"],
        ["저온 성능(-30°C)", "거의 유지", "20~40% 감소", "FCEV"],
        ["차량 가격", "7,000~8,000만", "4,000~6,000만", "BEV"],
        ["km당 에너지 비용", "100~120원", "40~60원", "BEV"],
        ["충전 인프라 수", "~1,200개소", "수백만 기", "BEV"],
    ]
    style_data_rows(table, data, size=12)

    # 에너지 효율 박스
    add_colored_box(slide, Inches(9.5), Inches(1.5), Inches(3.5), Inches(2.0),
                    NAVY, "에너지 효율 핵심\nFCEV: 최종 18~32%\nBEV: 최종 64~77%\n동일 재생에너지로\nBEV가 2~3배 더 많은\n주행거리 달성",
                    text_size=12)
    add_colored_box(slide, Inches(9.5), Inches(3.8), Inches(3.5), Inches(1.5),
                    GREEN, "FCEV 강점\n충전 3~5분\n주행거리 720km+\n저온 성능 유지",
                    text_size=14)
    add_colored_box(slide, Inches(9.5), Inches(5.5), Inches(3.5), Inches(1.0),
                    ACCENT_RED, "BEV 강점\n효율 2~3배 / 비용 절반\n인프라 압도적",
                    text_size=14)


def slide_21_tco_environment(prs):
    """TCO/환경성"""
    slide = setup_slide(prs, "TCO(총소유비용) & 환경성 비교", 21)

    # TCO 비교 테이블
    tf_label1 = add_body_textbox(slide, left=Inches(0.6), top=Inches(1.35),
                                  width=Inches(6.5), height=Inches(0.3))
    add_bullet(tf_label1, "승용차 TCO 비교 — 10년/15만km", first=True, size=15, bold=True, color=NAVY)

    tbl1 = slide.shapes.add_table(6, 4, Inches(0.6), Inches(1.7),
                                   Inches(6.5), Inches(2.5))
    t1 = tbl1.table
    t1.columns[0].width = Inches(2.0)
    t1.columns[1].width = Inches(1.5)
    t1.columns[2].width = Inches(1.5)
    t1.columns[3].width = Inches(1.5)
    style_header_row(t1, ["항목", "FCEV", "BEV", "내연기관"])
    data1 = [
        ["차량 구입(보조금 후)", "~3,950만", "~3,820만", "~3,500만"],
        ["연료/전기(10년)", "~1,800만", "~600만", "~2,000만"],
        ["유지보수(10년)", "~350만", "~230만", "~700만"],
        ["보험(10년)", "~380만", "~380만", "~400만"],
        ["10년 TCO", "~6,480만", "~4,650만", "~6,600만"],
    ]
    style_data_rows(t1, data1, size=12)

    # 환경성 비교 테이블
    tf_label2 = add_body_textbox(slide, left=Inches(0.6), top=Inches(4.4),
                                  width=Inches(6.5), height=Inches(0.3))
    add_bullet(tf_label2, "Well-to-Wheel CO₂ 배출량 (g/km)", first=True, size=15, bold=True, color=NAVY)

    tbl2 = slide.shapes.add_table(4, 4, Inches(0.6), Inches(4.7),
                                   Inches(6.5), Inches(1.8))
    t2 = tbl2.table
    t2.columns[0].width = Inches(2.5)
    t2.columns[1].width = Inches(1.5)
    t2.columns[2].width = Inches(1.3)
    t2.columns[3].width = Inches(1.2)
    style_header_row(t2, ["시나리오", "FCEV", "BEV", "내연기관"])
    data2 = [
        ["그레이수소+화석전기", "180~200", "100~150", "~200"],
        ["블루수소+평균전력망", "80~120", "50~80", "-"],
        ["그린수소+재생에너지", "0~10", "0~5", "-"],
    ]
    style_data_rows(t2, data2, size=12)

    # 오른쪽 핵심 박스
    add_colored_box(slide, Inches(7.5), Inches(1.5), Inches(5.3), Inches(1.5),
                    ACCENT_RED, "TCO 결론\nBEV가 FCEV 대비 약 28% 저렴\n연료비 격차가 핵심\n(600만 vs 1,800만 원/10년)",
                    text_size=14)
    add_colored_box(slide, Inches(7.5), Inches(3.3), Inches(5.3), Inches(1.5),
                    GREEN, "대형 트럭 — FCEV 역전 가능\nMcKinsey: 2030년 이후\n장거리 대형 트럭 TCO 동등/유리\n배터리 4~8톤→화물 20~30%↓",
                    text_size=14)
    add_colored_box(slide, Inches(7.5), Inches(5.1), Inches(5.3), Inches(1.3),
                    NAVY, "환경성 현실\n수소 99% 그레이수소\n현 상태에서 BEV가 유리\n그린수소 전환이 관건",
                    text_size=14)


def slide_22_optimal_applications(prs):
    """용도별 최적 기술"""
    slide = setup_slide(prs, "용도별 최적 기술 — BEV vs FCEV 역할 분담", 22)

    # 테이블
    tbl_shape = slide.shapes.add_table(10, 3, Inches(0.6), Inches(1.5),
                                       Inches(8.0), Inches(4.8))
    table = tbl_shape.table
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(4.0)
    style_header_row(table, ["용도", "최적 기술", "이유"])
    data = [
        ["도심 승용차", "BEV", "짧은 통근, 가정 충전, 비용 효율"],
        ["택시/라이드헤일링", "FCEV", "빠른 충전, 장시간 운행"],
        ["시내버스", "BEV/FCEV", "노선 길이에 따라 혼용"],
        ["장거리 고속버스", "FCEV", "주행거리, 충전 시간 우위"],
        ["소형 배송 트럭", "BEV", "도심 단거리, 야간 충전"],
        ["대형 장거리 트럭", "FCEV", "페이로드, 주행거리, 충전시간"],
        ["물류 지게차", "FCEV", "24시간 가동, 2분 충전, 냉동창고"],
        ["선박", "FCEV/수소", "장기 항해, 대용량 에너지"],
        ["항공기", "FCEV/e-Fuel", "에너지밀도 우위 필수"],
    ]
    style_data_rows(table, data, size=12)

    # 핵심 메시지 박스
    add_colored_box(slide, Inches(9.0), Inches(1.5), Inches(3.8), Inches(2.5),
                    NAVY, "BEV 영역\n도심 승용차\n소형 배송 트럭\n단거리·소형 차량\n→ BEV가 절대 우위",
                    text_size=14)
    add_colored_box(slide, Inches(9.0), Inches(4.3), Inches(3.8), Inches(2.5),
                    GREEN, "FCEV 영역\n대형 트럭 / 고속버스\n택시 / 지게차\n선박 / 항공\n→ FCEV가 비교우위",
                    text_size=14)


def slide_23_market_comparison(prs):
    """시장 전망 비교"""
    slide = setup_slide(prs, "BEV vs FCEV — 글로벌 판매 전망 비교", 23)

    # 판매 전망 테이블
    tf_label1 = add_body_textbox(slide, left=Inches(0.6), top=Inches(1.35),
                                  width=Inches(6.5), height=Inches(0.3))
    add_bullet(tf_label1, "글로벌 판매 전망", first=True, size=15, bold=True, color=NAVY)

    tbl1 = slide.shapes.add_table(4, 4, Inches(0.6), Inches(1.7),
                                   Inches(6.5), Inches(1.8))
    t1 = tbl1.table
    t1.columns[0].width = Inches(1.5)
    t1.columns[1].width = Inches(1.8)
    t1.columns[2].width = Inches(2.0)
    t1.columns[3].width = Inches(1.2)
    style_header_row(t1, ["구분", "2025년", "2030년", "2035년"])
    data1 = [
        ["BEV", "~1,700만 대", "~3,500~4,000만", "신차 60~75%"],
        ["FCEV", "~2만 대", "~20만 대", "~60만 대"],
        ["BEV/FCEV 비율", "850:1", "175~200:1", "100:1"],
    ]
    style_data_rows(t1, data1, size=13)

    # 기관별 전망 테이블
    tf_label2 = add_body_textbox(slide, left=Inches(0.6), top=Inches(3.7),
                                  width=Inches(6.5), height=Inches(0.3))
    add_bullet(tf_label2, "기관별 전망 요약", first=True, size=15, bold=True, color=NAVY)

    tbl2 = slide.shapes.add_table(5, 4, Inches(0.6), Inches(4.0),
                                   Inches(6.5), Inches(2.5))
    t2 = tbl2.table
    t2.columns[0].width = Inches(1.2)
    t2.columns[1].width = Inches(1.8)
    t2.columns[2].width = Inches(1.7)
    t2.columns[3].width = Inches(1.8)
    style_header_row(t2, ["기관", "BEV 전망", "FCEV 전망", "핵심 메시지"])
    data2 = [
        ["IEA", "EV의 95%+", "승용 0.2~0.5%", "수소는 비전기화 분야"],
        ["BNEF", "신차 42~58%", "신차 0.22%", "승용 FCEV 경쟁 불가"],
        ["McKinsey", "승용 주도", "트럭 15~25%", "상용차 TCO 경쟁적"],
        ["IRENA", "소형차 주류", "장거리 30%", "1.5°C에 수소 필수"],
    ]
    style_data_rows(t2, data2, size=11)

    # 제조사 전략 불릿
    tf = add_body_textbox(slide, left=Inches(7.5), top=Inches(1.5), width=Inches(5.3), height=Inches(5.0))
    add_bullet(tf, "제조사 전략 분류", first=True, size=16, bold=True, color=NAVY)
    add_bullet(tf, "", size=8)
    add_bullet(tf, "BEV 올인(수소 거부)", size=14, bold=True, color=ACCENT_RED)
    add_bullet(tf, "테슬라, BYD, 폭스바겐", level=1, size=13)
    add_bullet(tf, "", size=8)
    add_bullet(tf, "BEV 중심 + 수소 유보/철수", size=14, bold=True, color=ACCENT_ORANGE)
    add_bullet(tf, "GM, Ford, Stellantis(중단)", level=1, size=13)
    add_bullet(tf, "", size=8)
    add_bullet(tf, "BEV + FCEV 병행(가장 적극적)", size=14, bold=True, color=GREEN)
    add_bullet(tf, "현대차(넥쏘+XCIENT+HTWO)", level=1, size=13)
    add_bullet(tf, "도요타(미라이+BMW 협력)", level=1, size=13)
    add_bullet(tf, "BMW(2028 양산), 혼다(150kW FC)", level=1, size=13)


def slide_24_scenarios(prs):
    """공존/경쟁 시나리오"""
    slide = setup_slide(prs, "공존 vs 경쟁 시나리오 — BEV & FCEV의 미래", 24)

    # 시나리오 A
    add_colored_box(slide, Inches(0.6), Inches(1.5), Inches(3.9), Inches(5.0),
                    GREEN,
                    "시나리오 A: 공존(보완 기술)\n확률 60~70%\n\n"
                    "가장 유력한 시나리오\n\n"
                    "• 도심/단거리 승용:\n  BEV 95%+ 주류\n\n"
                    "• 장거리/대형/고사이클:\n  FCEV 15~25%\n\n"
                    "• 수소: 산업/발전/저장\n  역할 분담",
                    text_size=13)

    # 시나리오 B
    add_colored_box(slide, Inches(4.8), Inches(1.5), Inches(3.9), Inches(5.0),
                    ACCENT_ORANGE,
                    "시나리오 B: BEV 지배\n확률 20~30%\n\n"
                    "FCEV 주변화 가능\n\n"
                    "• 전고체 배터리 상용화\n  (2027~2030)\n\n"
                    "• 초급속 충전 보급\n\n"
                    "• 수소 인프라 투자 부족\n\n"
                    "• 선박/항공/산업 외\n  거의 소멸",
                    text_size=13)

    # 시나리오 C
    add_colored_box(slide, Inches(9.0), Inches(1.5), Inches(3.8), Inches(5.0),
                    ACCENT_RED,
                    "시나리오 C: FCEV 급성장\n확률 5~10%\n\n"
                    "현실 가능성 낮음\n\n"
                    "• 그린수소 $1/kg 달성\n\n"
                    "• 연료전지 $30/kW 이하\n\n"
                    "• 대규모 인프라 투자\n\n"
                    "• 전체 신차 5~10%\n  확대 가능",
                    text_size=13)


def slide_25_five_variables(prs):
    """5대 핵심 변수"""
    slide = setup_slide(prs, "FCEV 성패를 결정할 5대 핵심 변수", 25)

    # 5대 변수 테이블
    tbl1 = slide.shapes.add_table(6, 4, Inches(0.6), Inches(1.5),
                                   Inches(12.1), Inches(3.0))
    t1 = tbl1.table
    t1.columns[0].width = Inches(2.5)
    t1.columns[1].width = Inches(2.5)
    t1.columns[2].width = Inches(2.5)
    t1.columns[3].width = Inches(4.6)
    style_header_row(t1, ["변수", "현재", "목표", "달성 시 영향"])
    data1 = [
        ["그린수소 비용", "$3~8/kg", "$1~2/kg", "연료비 경쟁력 확보"],
        ["연료전지 비용", "$80~100/kW", "$35/kW", "차량 가격 내연기관 수준"],
        ["충전 인프라", "~1,200개소", "수만 개소", "소비자 접근성 해소"],
        ["스택 내구성", "5,000~8,000시간", "30,000시간", "상용차 TCO 확보"],
        ["정부 정책 일관성", "불확실", "장기 안정적", "민간 투자 유치"],
    ]
    style_data_rows(t1, data1, size=13)

    # 전문가 합의 테이블
    tf_label = add_body_textbox(slide, left=Inches(0.6), top=Inches(4.7),
                                 width=Inches(12.0), height=Inches(0.3))
    add_bullet(tf_label, "전문가 합의 사항", first=True, size=16, bold=True, color=NAVY)

    tbl2 = slide.shapes.add_table(6, 2, Inches(0.6), Inches(5.0),
                                   Inches(12.1), Inches(1.8))
    t2 = tbl2.table
    t2.columns[0].width = Inches(0.5)
    t2.columns[1].width = Inches(11.6)
    style_header_row(t2, ["#", "합의 내용"])
    data2 = [
        ["1", "승용차 시장은 BEV의 압도적 승리가 거의 확실"],
        ["2", "대형 상용차(트럭/버스)에서 FCEV의 역할이 존재"],
        ["3", "수소는 모빌리티를 넘어 산업/발전/저장 분야에서 필수적"],
        ["4", "2050 탄소중립 달성에 수소와 BEV 모두 필요"],
        ["5", "FCEV의 성패는 수소 비용 하락 속도에 달려 있음"],
    ]
    style_data_rows(t2, data2, size=12)


def slide_26_korea_strategy(prs):
    """한국 전략 권고"""
    slide = setup_slide(prs, "한국에 대한 전략적 권고 — 5대 핵심 방향", 26)

    strategies = [
        ("1. 승용 FCEV\n보조금 점진적 축소",
         "넥쏘 42.9% 점유율은 보조금 기반\n장기 지속가능하지 않음\n보조금 의존 구조에서 탈피 필요",
         ACCENT_RED),
        ("2. 대형 상용차\n선택과 집중",
         "XCIENT 유럽/북미 2,000만km 실적\n수소 트럭/버스 공격적 확대\n글로벌 선두 + BEV 대비 비교우위",
         GREEN),
        ("3. 핵심 소재\n국산화 가속",
         "전해질막 10~20%, 촉매 10~15%\n해외 의존도 → 공급망 리스크\n선도국 대비 5~7년 격차 해소",
         NAVY),
        ("4. 그린수소\n전환 가속",
         "수소 90%가 그레이수소\n'친환경' 명분 성립 안 됨\n수전해 기반 비중 확대 시급",
         ACCENT_ORANGE),
        ("5. K-조선·선박\n시너지",
         "한국 조선 산업 + 수소 FC 결합\n수소 선박 = 독자적 블루오션\n한국만의 강점 발휘 영역",
         LIGHT_NAVY),
    ]

    for i, (title, desc, color) in enumerate(strategies):
        y = Inches(1.5) + Inches(i * 1.1)
        add_colored_box(slide, Inches(0.6), y, Inches(2.8), Inches(0.95),
                        color, title, text_size=13, bold=True)
        txBox = slide.shapes.add_textbox(Inches(3.7), y, Inches(9.0), Inches(0.95))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.05)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = desc
        set_font(run, size=14, color=DARK_GRAY)


def slide_27_key_numbers(prs):
    """핵심 수치 대시보드"""
    slide = setup_slide(prs, "핵심 수치 대시보드 — 한눈에 보는 수소차 시장", 27)

    items = [
        ("글로벌 FCEV 판매\n2024년", "~16,000대\n(BEV의 0.1%)", NAVY),
        ("넥쏘 글로벌 점유율\n2025년", "42.9%\n세계 1위", GREEN),
        ("BEV vs FCEV\n규모 격차", "약 850배\n(1,380만 vs 1.6만)", ACCENT_RED),
        ("수소 충전소\n전 세계", "~1,200개\n(EV 수백만 기)", LIGHT_NAVY),
        ("충전소 설치비\n수소 vs EV", "30~100배\n비싸", ACCENT_ORANGE),
        ("TCO 격차\n승용 10년", "BEV 28%\n더 저렴", ACCENT_BLUE),
        ("5대 그룹\n총 투자", "~43.4조 원\n(2030년까지)", NAVY),
        ("공존 시나리오\n확률", "60~70%\n가장 유력", GREEN),
    ]

    for i, (title, value, color) in enumerate(items):
        col = i % 4
        row = i // 4
        x = Inches(0.6 + col * 3.15)
        y = Inches(1.5 + row * 2.8)
        add_colored_box(slide, x, y, Inches(2.9), Inches(2.5),
                        color, f"{title}\n\n{value}", text_size=14)


def slide_28_conclusion(prs):
    """종합 결론"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, NAVY)

    # 상단 녹색 라인
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Inches(0.6), SLIDE_WIDTH, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()

    # 제목
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "종합 결론"
    set_font(run, size=40, bold=True, color=WHITE)

    # 핵심 요약 테이블 형태 (색상 박스)
    items = [
        ("승용차", "BEV의 확정적 승자 — FCEV는 0.5% 미만 니치 시장"),
        ("대형 상용차", "FCEV의 비교우위 영역 — 트럭/버스/선박에서 역할 존재"),
        ("에너지 시스템", "수소는 산업/발전/저장 분야에서 필수적 — 모빌리티 넘어 확장"),
    ]

    for i, (title, desc) in enumerate(items):
        y = Inches(1.9) + Inches(i * 0.85)
        add_colored_box(slide, Inches(0.8), y, Inches(2.8), Inches(0.7),
                        GREEN, title, text_size=16, bold=True)
        txB = slide.shapes.add_textbox(Inches(4.0), y, Inches(8.5), Inches(0.7))
        tfB = txB.text_frame
        tfB.word_wrap = True
        pB = tfB.paragraphs[0]
        pB.alignment = PP_ALIGN.LEFT
        runB = pB.add_run()
        runB.text = desc
        set_font(runB, size=17, color=WHITE)

    # 인용문
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(11.7), Inches(1.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = '"수소는 만능이 아니다. 그러나 수소 없이 탄소중립은 불가능하다."'
    set_font(run2, size=22, bold=True, color=GREEN)

    p3 = tf2.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "BEV와 FCEV는 경쟁이 아닌 공존의 관계이며, 각 기술이 최적인 영역에서 역할을 분담하는 것이\n2050 탄소중립을 향한 가장 현실적인 경로이다."
    set_font(run3, size=16, color=RGBColor(0xBB, 0xD5, 0xED))

    # 한국 전략 5대 방향
    txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.4))
    tf3 = txBox3.text_frame
    p4 = tf3.paragraphs[0]
    p4.alignment = PP_ALIGN.LEFT
    run4 = p4.add_run()
    run4.text = "한국 전략: 승용 보조금 축소 → 대형 상용차 집중 → 핵심소재 국산화 → 그린수소 전환 → K-조선 시너지"
    set_font(run4, size=15, bold=True, color=ACCENT_ORANGE)

    # 하단 녹색 라인
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Inches(6.8), SLIDE_WIDTH, Inches(0.06)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = GREEN
    shape2.line.fill.background()

    add_slide_number(slide, TOTAL_SLIDES)


# ── 메인 ──────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_01_cover(prs)
    slide_02_toc(prs)
    slide_03_global_market(prs)
    slide_04_country_status(prs)
    slide_05_key_models(prs)
    slide_06_other_models(prs)
    slide_07_charging_infra(prs)
    slide_08_market_outlook(prs)
    slide_09_government_policy(prs)
    slide_10_commercial_vehicles(prs)
    slide_11_other_mobility(prs)
    slide_12_stack_technology(prs)
    slide_13_cost_storage(prs)
    slide_14_korea_market(prs)
    slide_15_korea_infra_crisis(prs)
    slide_16_korea_policy(prs)
    slide_17_korea_investment(prs)
    slide_18_korea_supply_chain(prs)
    slide_19_korea_global_status(prs)
    slide_20_bev_vs_fcev_specs(prs)
    slide_21_tco_environment(prs)
    slide_22_optimal_applications(prs)
    slide_23_market_comparison(prs)
    slide_24_scenarios(prs)
    slide_25_five_variables(prs)
    slide_26_korea_strategy(prs)
    slide_27_key_numbers(prs)
    slide_28_conclusion(prs)

    prs.save(OUTPUT_PATH)
    print(f"PPT 생성 완료: {OUTPUT_PATH}")
    print(f"총 {TOTAL_SLIDES}장 슬라이드")


if __name__ == "__main__":
    main()
