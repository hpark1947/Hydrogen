"""
연료전지 발표 PPT 생성 스크립트
- 25장 슬라이드, 비즈니스 프로페셔널 디자인
- python-pptx 기반
- 글로벌 동향, 미국·중국, 한국 현황 및 전망
- 로드맵 달성률, 부품 국산화/기술격차, 충전소 수익성 위기 추가
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 상수 ──────────────────────────────────────────────
NAVY = RGBColor(0x1B, 0x3A, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_NAVY = RGBColor(0x2C, 0x5F, 0x8A)
ACCENT_BLUE = RGBColor(0x34, 0x98, 0xDB)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_ORANGE = RGBColor(0xF3, 0x9C, 0x12)
TABLE_HEADER_BG = RGBColor(0x1B, 0x3A, 0x5C)
TABLE_ROW_ALT = RGBColor(0xE8, 0xF0, 0xF8)
TABLE_ROW_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_NAME = "맑은 고딕"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
TOTAL_SLIDES = 25

OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "연료전지_발표자료.pptx")


# ── 유틸리티 함수 ─────────────────────────────────────
def set_font(run, size=18, bold=False, color=DARK_GRAY, name=FONT_NAME):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_background(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_navy_header_bar(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_WIDTH, Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()


def add_green_accent_line(slide, top=Inches(1.2)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), top, SLIDE_WIDTH, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()


def add_slide_number(slide, num):
    txBox = slide.shapes.add_textbox(
        Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{num} / {TOTAL_SLIDES}"
    set_font(run, size=12, color=MEDIUM_GRAY)


def add_footer_line(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(7.1), Inches(12.333), Inches(0.02)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()


def add_title_text(slide, title_text, left=Inches(0.6), top=Inches(0.2),
                   width=Inches(12), height=Inches(0.9)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    set_font(run, size=34, bold=True, color=WHITE)
    return txBox


def add_body_textbox(slide, left=Inches(0.8), top=Inches(1.6),
                     width=Inches(11.7), height=Inches(5.2)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_bullet(tf, text, level=0, size=16, bold=False, color=DARK_GRAY,
               space_after=Pt(6), first=False):
    if first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.level = level
    p.space_after = space_after
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    return p


def set_cell(table, row, col, text, size=12, bold=False, color=DARK_GRAY,
             alignment=PP_ALIGN.LEFT, fill_color=None):
    cell = table.cell(row, col)
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color


def style_header_row(table, headers, size=14):
    for i, h in enumerate(headers):
        set_cell(table, 0, i, h, size=size, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER, fill_color=TABLE_HEADER_BG)


def style_data_rows(table, data, start_row=1, size=13):
    for r_idx, row_data in enumerate(data):
        row_num = start_row + r_idx
        bg = TABLE_ROW_ALT if row_num % 2 == 0 else TABLE_ROW_WHITE
        for c_idx, val in enumerate(row_data):
            set_cell(table, row_num, c_idx, val, size=size, fill_color=bg)


def setup_slide(prs, title_text, slide_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_navy_header_bar(slide)
    add_green_accent_line(slide)
    add_title_text(slide, title_text)
    add_footer_line(slide)
    add_slide_number(slide, slide_num)
    return slide


def add_colored_box(slide, left, top, width, height, fill_color, text,
                    text_size=13, text_color=WHITE, bold=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_font(run, size=text_size, bold=bold, color=text_color)
    return shape


# ── 슬라이드 생성 ─────────────────────────────────────

def slide_01_cover(prs):
    """표지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, NAVY)

    # 상단 녹색 라인
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Inches(2.0), SLIDE_WIDTH, Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()

    # 제목
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.4), Inches(10.3), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "연료전지(Fuel Cell) 산업 동향"
    set_font(run, size=46, bold=True, color=WHITE)

    # 부제
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.9), Inches(10.3), Inches(1.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "글로벌 시장 동향 | 미국·중국 경쟁 구도 | 한국의 현주소와 전망"
    set_font(run2, size=22, bold=False, color=RGBColor(0xBB, 0xD5, 0xED))

    # 날짜
    txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.6))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "2026. 02"
    set_font(run3, size=18, color=RGBColor(0x88, 0xAA, 0xCC))

    # 하단 녹색 라인
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Inches(5.7), SLIDE_WIDTH, Inches(0.08)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = GREEN
    shape2.line.fill.background()


def slide_02_toc(prs):
    """목차"""
    slide = setup_slide(prs, "목차 (Table of Contents)", 2)
    sections = [
        ("Part 1", "글로벌 연료전지 시장 동향", "시장 규모, 유형별 동향, 응용 분야, 주요국 정책"),
        ("Part 2", "미국 연료전지 산업", "정책, 주요 기업, R&D, 응용 분야"),
        ("Part 3", "중국 연료전지 산업", "정책, 주요 기업, 기술 수준, 상용차 중심"),
        ("Part 4", "미국 vs 중국 경쟁 구도", "전략 비교, 강점·약점, 향후 전망"),
        ("Part 5", "한국의 현주소와 전망", "시장, 기업, R&D, 인프라, 달성률, 기술격차, 충전소위기"),
    ]
    colors = [NAVY, LIGHT_NAVY, ACCENT_BLUE, ACCENT_ORANGE, GREEN]

    for i, (part, title, desc) in enumerate(sections):
        y = Inches(1.6) + Inches(i * 1.05)
        add_colored_box(slide, Inches(0.8), y, Inches(1.5), Inches(0.8),
                        colors[i], part, text_size=16, bold=True)
        txBox = slide.shapes.add_textbox(Inches(2.6), y, Inches(9.5), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        set_font(run, size=20, bold=True, color=DARK_GRAY)
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = desc
        set_font(run2, size=14, color=MEDIUM_GRAY)


def slide_03_global_market(prs):
    """글로벌 시장 규모"""
    slide = setup_slide(prs, "글로벌 연료전지 시장 규모 및 성장 전망", 3)

    tf = add_body_textbox(slide, top=Inches(1.5), height=Inches(1.2))
    add_bullet(tf, "2025년 시장 규모: 약 107.6억~129.4억 달러", first=True, size=18, bold=True)
    add_bullet(tf, "2030년까지 CAGR 20~27%로 고속 성장 전망", size=18, bold=True)
    add_bullet(tf, "아시아태평양 지역이 시장 점유율 42.7%로 최대 시장", size=16)

    # 표: 기관별 전망
    tbl_shape = slide.shapes.add_table(5, 3, Inches(0.8), Inches(3.2), Inches(7.5), Inches(3.0))
    table = tbl_shape.table
    table.columns[0].width = Inches(3.0)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.0)
    style_header_row(table, ["조사기관", "2030년 전망", "CAGR"])
    data = [
        ["Grand View Research", "369.8억 달러", "27.1%"],
        ["MarketsandMarkets", "181.6억 달러", "26.3%"],
        ["Wissen Research", "210억 달러", "25.0%"],
        ["MarkNtel Advisors", "81.3억 달러", "20.4%"],
    ]
    style_data_rows(table, data)

    # 오른쪽 핵심 박스
    add_colored_box(slide, Inches(9.0), Inches(3.2), Inches(3.8), Inches(1.2),
                    NAVY, "2034년 장기 전망\n372억~957억 달러", text_size=16)
    add_colored_box(slide, Inches(9.0), Inches(4.6), Inches(3.8), Inches(1.2),
                    GREEN, "핵심 성장 동인\nAI 데이터센터 + 상용차", text_size=16)


def slide_04_fc_types(prs):
    """연료전지 유형별 동향"""
    slide = setup_slide(prs, "연료전지 유형별 동향", 4)

    tbl_shape = slide.shapes.add_table(5, 4, Inches(0.6), Inches(1.5),
                                       Inches(12.1), Inches(5.0))
    table = tbl_shape.table
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(4.5)
    table.columns[3].width = Inches(3.8)

    style_header_row(table, ["유형", "시장 점유율", "주요 특성", "핵심 동향"])
    data = [
        ["PEMFC\n(고분자전해질막)", "60~67%\n(1위)",
         "작동온도 60~80°C\n빠른 기동, 높은 출력밀도",
         "백금 사용량 90% 감소 달성\nFe-N-C 촉매 0.85 W/cm²"],
        ["SOFC\n(고체산화물)", "~20%",
         "작동온도 600~1,000°C\n효율 60~75%, 다연료 가능",
         "AI 데이터센터 전력 공급 폭발적 성장\nBloom Energy 1GW 공급계약"],
        ["MCFC\n(용융탄산염)", "~15%",
         "작동온도 ~650°C\n대규모 산업/발전용",
         "FuelCell Energy 주도\n대형 분산발전"],
        ["PAFC\n(인산형)", "기타",
         "중간 온도 작동\n도시가스 기반 분산발전",
         "두산퓨얼셀 PureCell 시리즈\n안정적 발전용"],
    ]
    style_data_rows(table, data, size=12)


def slide_05_applications(prs):
    """응용 분야별 동향"""
    slide = setup_slide(prs, "주요 응용 분야별 동향", 5)

    # 고정형 발전
    add_colored_box(slide, Inches(0.6), Inches(1.5), Inches(3.7), Inches(5.0),
                    NAVY,
                    "고정형 발전\n(시장 점유율 68.8%)\n\n"
                    "━━━━━━━━━━━━\n\n"
                    "AI 데이터센터 전력 수요\n급증이 핵심 성장 동인\n\n"
                    "2026년 = 'SOFC의 해'\n\n"
                    "2030년까지 데이터센터\n38%가 온사이트 발전 활용",
                    text_size=14, bold=False)

    # 수송
    add_colored_box(slide, Inches(4.6), Inches(1.5), Inches(3.7), Inches(5.0),
                    LIGHT_NAVY,
                    "수송 분야\n(최고 성장률)\n\n"
                    "━━━━━━━━━━━━\n\n"
                    "2025 수소차 16,011대 판매\n(전년 대비 +24.4%)\n\n"
                    "현대차 42.9% 세계 1위\n\n"
                    "수소 상용차 CAGR 47.7%\n(2025~2032)",
                    text_size=14, bold=False)

    # 신규 응용
    add_colored_box(slide, Inches(8.6), Inches(1.5), Inches(4.1), Inches(5.0),
                    ACCENT_BLUE,
                    "신규 응용 분야\n\n"
                    "━━━━━━━━━━━━\n\n"
                    "[해양] 한화에어로 200kW\n선박용 연료전지, DNV 인증\n\n"
                    "[항공] ZeroAvia\n1.0 kW/kg 출력밀도 달성\n\n"
                    "[군사] UAV 비행시간\n2h → 6h+ 연장\n\n"
                    "[휴대용] CAGR 14%+",
                    text_size=14, bold=False)


def slide_06_global_policy(prs):
    """주요국 정책"""
    slide = setup_slide(prs, "세계 주요국 수소/연료전지 정책 현황", 6)

    tbl_shape = slide.shapes.add_table(6, 3, Inches(0.6), Inches(1.5),
                                       Inches(12.1), Inches(5.2))
    table = tbl_shape.table
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(3.5)
    table.columns[2].width = Inches(7.1)

    style_header_row(table, ["국가", "핵심 정책", "주요 내용"])
    data = [
        ["미국", "Hydrogen Shot + IRA",
         "2031년까지 $1/kg 목표 | H2Hubs 70억$ (일부 삭감) | 45V 세액공제"],
        ["EU", "EU 수소전략 + REPowerEU",
         "2030년 재생수소 2,000만톤 | 전해조 40GW | 1.845억 유로 공모"],
        ["일본", "수소기본전략 (2023 개정)",
         "관민 15조 엔(1,000억$) | 에네팜 50만대+ | 세계 최초 수소전략(2017)"],
        ["한국", "수소경제 로드맵 + 수소법",
         "세계 최초 수소법(2020) | 2040년 FCEV 620만대 | 발전용 15GW"],
        ["중국", "수소 중장기 계획 (2021-2035)",
         "2025년 FCV 5만대 목표 | 에너지법에 수소 포함 | 24개 성시 33건 정책"],
    ]
    style_data_rows(table, data, size=13)


def slide_07_us_policy(prs):
    """미국 연료전지 정책"""
    slide = setup_slide(prs, "[미국] 수소/연료전지 정책 상세", 7)
    tf = add_body_textbox(slide, top=Inches(1.5))

    add_bullet(tf, "국가 청정 수소 전략 및 로드맵 (2023)", first=True, size=18, bold=True, color=NAVY)
    add_bullet(tf, "2030년까지 연간 수소 생산 1,000만 톤 목표 (현재 7~9 MMTpa 전망)", level=1, size=15)

    add_bullet(tf, "인플레이션 감축법(IRA) - Section 45V", size=18, bold=True, color=NAVY)
    add_bullet(tf, "수소 1kg당 최대 $3 세액공제 | 트럼프 행정부 하 공제기간 10년→2년 단축", level=1, size=15)

    add_bullet(tf, "지역 청정 수소 허브 (H2Hubs)", size=18, bold=True, color=NAVY)
    add_bullet(tf, "7개 허브에 70억 달러 투자 → 트럼프 행정부, 서부 2개 허브 22억$ 취소", level=1, size=15)
    add_bullet(tf, "캘리포니아 ARCHES 12억$ 취소, 태평양 북서부 10억$ 삭감", level=1, size=15)

    add_bullet(tf, "Hydrogen Shot (DOE)", size=18, bold=True, color=NAVY)
    add_bullet(tf, "2031년까지 청정수소 비용 80% 절감, $1/kg 달성 목표", level=1, size=15)
    add_bullet(tf, "PEM 전해조 자본비용 80% 절감(2005~), 연료전지 시스템 70% 절감(2008~)", level=1, size=15)

    add_bullet(tf, "정책 리스크: 행정부 교체에 따른 불확실성이 투자 저해 요인", size=16, bold=True,
               color=ACCENT_RED)


def slide_08_us_companies(prs):
    """미국 주요 기업"""
    slide = setup_slide(prs, "[미국] 주요 연료전지 기업", 8)

    tbl_shape = slide.shapes.add_table(6, 4, Inches(0.6), Inches(1.5),
                                       Inches(12.1), Inches(5.2))
    table = tbl_shape.table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(4.0)
    table.columns[3].width = Inches(3.4)

    style_header_row(table, ["기업", "핵심 기술", "주요 실적", "특이사항"])
    data = [
        ["Bloom Energy", "SOFC",
         "2025 매출 20.2억$ (+37.3%)\n수주잔고 200억$",
         "업계 유일 수익 기업\nAI 데이터센터 핵심"],
        ["Plug Power", "PEMFC + 전해조",
         "2025 9개월 4.847억$\nGenEco 230MW+ 파이프라인",
         "Amazon 지게차 15,000대+\n그린수소 통합 솔루션"],
        ["Ballard Power", "PEMFC 스택",
         "Q3 매출 3,250만$ (+120%)\n수주잔고 1.328억$",
         "캐나다 본사\nWeichai와 중국 JV"],
        ["Cummins/Accelera", "연료전지 + 전해조",
         "린데에 35MW 전해조 납품\n전체 매출 337억$",
         "H2-ICE 병행 추진\n4.58억$ 구조조정"],
        ["Nikola", "수소 트럭", "2025.2 파산 신청\n트럭 95대 리콜",
         "Chapter 11\n수소트럭 상용화 어려움"],
    ]
    style_data_rows(table, data, size=11)


def slide_09_us_applications(prs):
    """미국 응용 분야"""
    slide = setup_slide(prs, "[미국] 연료전지 응용 분야", 9)
    tf = add_body_textbox(slide, top=Inches(1.5))

    add_bullet(tf, "물류/지게차 — 가장 상용화된 분야", first=True, size=18, bold=True, color=NAVY)
    add_bullet(tf, "Amazon 15,000대+ 운영, 5,000대 추가 계획 | Walmart 대규모 도입", level=1, size=15)
    add_bullet(tf, "3분 충전, 교대근무 시 일관된 출력, 배터리 교체 불필요", level=1, size=15)

    add_bullet(tf, "AI 데이터센터 분산전원 — 최대 성장 분야", size=18, bold=True, color=NAVY)
    add_bullet(tf, "Bloom Energy: AEP 1GW 계약, Brookfield 50억$ 파트너십", level=1, size=15)
    add_bullet(tf, "전력망 연결 수년 소요 → SOFC로 신속한 현장 배치 가능", level=1, size=15)

    add_bullet(tf, "군사 응용 — 독보적 영역", size=18, bold=True, color=NAVY)
    add_bullet(tf, "시장규모: 2024년 12억$ → 2033년 39억$ (CAGR 13.8%)", level=1, size=15)
    add_bullet(tf, "보병 휴대전원(무게 50% 절감), UAV(비행 6h+), 잠수함(무소음)", level=1, size=15)

    add_bullet(tf, "수소충전소: 전국 ~51개 (95%가 캘리포니아) — 인프라 극히 부족",
               size=16, bold=True, color=ACCENT_RED)


def slide_10_china_market(prs):
    """중국 시장"""
    slide = setup_slide(prs, "[중국] 연료전지 시장 현황", 10)

    # 좌측: 핵심 지표 표
    tbl_shape = slide.shapes.add_table(5, 3, Inches(0.6), Inches(1.5),
                                       Inches(7.0), Inches(2.8))
    table = tbl_shape.table
    table.columns[0].width = Inches(2.8)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(2.2)
    style_header_row(table, ["항목", "목표", "실적/달성률"])
    data = [
        ["연료전지차(FCV)", "50,000대", "~30,000대 (60%)"],
        ["수소충전소", "1,200개소", "~540개소 (45%)"],
        ["그린수소 생산", "10~20만 톤/년", "진행 중"],
        ["핵심 부품 국산화율", "-", "70% 달성"],
    ]
    style_data_rows(table, data)

    # 우측: 핵심 포인트 박스
    add_colored_box(slide, Inches(8.2), Inches(1.5), Inches(4.5), Inches(1.2),
                    NAVY, "시장 규모\n2023년 15.4억$ → 2026년 35.5억$", text_size=15)
    add_colored_box(slide, Inches(8.2), Inches(2.9), Inches(4.5), Inches(1.2),
                    GREEN, "국산 스택\n수입품 대비 60% 저렴", text_size=15)

    # 하단: 정책 요약
    tf = add_body_textbox(slide, top=Inches(4.6), height=Inches(2.0))
    add_bullet(tf, "중국 수소 에너지 중장기 계획 (2021~2035) — NDRC 발표", first=True, size=17, bold=True, color=NAVY)
    add_bullet(tf, "3단계 로드맵: 1단계(~2025) 프레임워크 → 2단계(~2030) 혁신 시스템 → 3단계(~2035) 다양한 생태계", level=1, size=14)
    add_bullet(tf, "2024.11: 24개 성/시에서 33개 신규 수소정책 | 2025.4: 에너지법에 수소 포함", level=1, size=14)
    add_bullet(tf, "2025년 재무부 FCV 보조금: 3.21억$ | 5개 시범 도시 클러스터 운영", level=1, size=14)


def slide_11_china_companies(prs):
    """중국 주요 기업"""
    slide = setup_slide(prs, "[중국] 주요 연료전지 기업", 11)

    tbl_shape = slide.shapes.add_table(6, 4, Inches(0.6), Inches(1.5),
                                       Inches(12.1), Inches(5.2))
    table = tbl_shape.table
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(4.0)
    table.columns[3].width = Inches(3.6)

    style_header_row(table, ["기업", "본사", "핵심 사업", "특이사항"])
    data = [
        ["SinoHytec (이화통)", "베이징",
         "수소 연료전지 엔진 개발\nIPO 14억 위안(2.13억$)",
         "Toyota 등 5개사와 공동개발\n시스템 시장 선두"],
        ["Refire (리파이어)", "상하이",
         "연료전지 엔진 제조\n(트럭·버스용)",
         "생산능력 1,000대→20,000대\n선제적 확대"],
        ["Weichai Power", "산둥성",
         "Ballard와 JV (51:49)\nMEA 장기 독점 공급",
         "SOFC 전략 병행\n국가급 핵심 멤버"],
        ["SAIC Motor", "상하이",
         "수소차 10종 출시 계획\nMaxus EUNIQ 7",
         "중국 최대 자동차 그룹\n수소차 전략 추진"],
        ["Sinosynergy", "-",
         "CCM 1.5만개/일, MEA 1만개/일\n핵심부품 자체 개발",
         "국산 스택 60% 저렴\n고객 70개사+"],
    ]
    style_data_rows(table, data, size=11)


def slide_12_china_applications(prs):
    """중국 상용차 중심 응용"""
    slide = setup_slide(prs, "[중국] 상용차 중심 대규모 보급", 12)
    tf = add_body_textbox(slide, top=Inches(1.5))

    add_bullet(tf, "중국 연료전지 시장의 60%+ = 수송 부문, 상용차가 압도적 비중",
               first=True, size=17, bold=True, color=NAVY)

    add_bullet(tf, "2025년 12월 주요 배치 사례", size=17, bold=True, color=NAVY)
    add_bullet(tf, "광저우 버스그룹: 수소 시내버스 450대 입찰 (Skywell 250대 + King Long 200대)", level=1, size=15)
    add_bullet(tf, "저장 롄허수소: 연료전지 세미트럭/덤프트럭 200대 공급", level=1, size=15)
    add_bullet(tf, "HTWO 광저우: 18톤 수소 중형트럭 200대 납품", level=1, size=15)
    add_bullet(tf, "친링모터스: 신장 하미시, 세계 최대 수소 상용차 시범 프로젝트", level=1, size=15)

    add_bullet(tf, "수소충전소 인프라", size=17, bold=True, color=NAVY)
    add_bullet(tf, "2024년 누적 540개소 (세계 최다) | 광둥 68개, 산둥/저장 각 20개+", level=1, size=15)
    add_bullet(tf, "Sinopec: 1,000개소 건설 목표 | 13개 성/시가 10개+ 보유", level=1, size=15)

    add_bullet(tf, "화물 운송 회랑", size=17, bold=True, color=NAVY)
    add_bullet(tf, "20개+ 주요 회랑에서 200대+ 운행, 40개+ 충전소 지원", level=1, size=15)


def slide_13_us_vs_china(prs):
    """미중 비교"""
    slide = setup_slide(prs, "미국 vs 중국 연료전지 경쟁 구도", 13)

    tbl_shape = slide.shapes.add_table(8, 3, Inches(0.6), Inches(1.5),
                                       Inches(12.1), Inches(5.2))
    table = tbl_shape.table
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(4.8)
    table.columns[2].width = Inches(4.8)

    style_header_row(table, ["비교 항목", "미국", "중국"])
    data = [
        ["전략 방향", "원천기술 + 고부가가치\n(데이터센터, 군사)", "대규모 보급 + 비용절감\n(상용차 중심)"],
        ["수소충전소", "~51개 (95% 캘리포니아)", "~540개 (세계 최다)"],
        ["FCV 보유", "소규모 (캘리포니아 중심)", "~30,000대"],
        ["핵심 기술력", "SOFC/PEMFC 원천기술\n촉매/MEA 선도", "핵심부품 국산화 70%\n비용 절감 우위"],
        ["대표 기업", "Bloom Energy, Plug Power", "SinoHytec, Refire, Weichai"],
        ["주요 강점", "원천기술, 군사, AI DC", "대규모 보급, 인프라, 정책 일관성"],
        ["주요 약점", "인프라 부족, 정책 불확실성", "기술 내구성, 핵심소재 수입"],
    ]
    style_data_rows(table, data, size=12)


def slide_14_korea_market(prs):
    """한국 시장 현황"""
    slide = setup_slide(prs, "[한국] 연료전지 시장 현황", 14)

    # 핵심 지표 박스 4개
    boxes = [
        ("발전용 연료전지\n1,036MW\n세계 최초 1GW 돌파", NAVY),
        ("수소차 등록\n4만 대 돌파\n넥쏘 점유율 42.9%", LIGHT_NAVY),
        ("5대 그룹 투자\n43.4조 원\n밸류체인 전 단계", ACCENT_BLUE),
        ("건물용 시장\n2030년\n3.86조 원", GREEN),
    ]
    for i, (text, color) in enumerate(boxes):
        add_colored_box(slide, Inches(0.6 + i * 3.1), Inches(1.5),
                        Inches(2.8), Inches(1.8), color, text, text_size=16)

    # 하단 상세 정보
    tf = add_body_textbox(slide, top=Inches(3.8), height=Inches(2.8))
    add_bullet(tf, "발전용 연료전지", first=True, size=17, bold=True, color=NAVY)
    add_bullet(tf, "세계 최초 1GW 돌파 | 2021년 142.3MW 설치, 세계 시장 45% 점유", level=1, size=15)

    add_bullet(tf, "수소차", size=17, bold=True, color=NAVY)
    add_bullet(tf, "현대 넥쏘 2025년 6,861대 판매 (+78.9%) | 신형 넥쏘 월 1,000대+ 출고", level=1, size=15)

    add_bullet(tf, "수소버스", size=17, bold=True, color=NAVY)
    add_bullet(tf, "2024년 1,000대+ 신규 보급 (전년 대비 277% 급증) | 2025년 목표 2,000대", level=1, size=15)


def slide_15_korea_policy(prs):
    """한국 정책"""
    slide = setup_slide(prs, "[한국] 수소경제 정책 체계", 15)

    tbl_shape = slide.shapes.add_table(7, 3, Inches(0.6), Inches(1.5),
                                       Inches(12.1), Inches(5.2))
    table = tbl_shape.table
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(6.6)

    style_header_row(table, ["정책명", "시기", "주요 내용"])
    data = [
        ["수소경제 활성화 로드맵", "2019.1", "2040년 FCEV 620만대, 충전소 1,200기, 발전 15GW"],
        ["수소법 (세계 최초)", "2020.2", "수소경제 육성 및 수소 안전관리 법률 제정"],
        ["수소경제 이행 기본계획", "2021", "2030 수소 390만톤, 2050 2,790만톤 공급 목표"],
        ["청정수소 인증제", "2024", "세계 최초 청정수소발전 입찰시장 개설 (낙찰률 11.5%)"],
        ["CHPS (의무화제도)", "2024~", "RPS에서 분리, 별도 의무시장 | 2025년 입찰 전격 취소"],
        ["전력수급기본계획", "2027~", "연간 200MW(~2030) → 150MW(~2036) → 100MW 증설"],
    ]
    style_data_rows(table, data, size=12)


def slide_16_korea_companies(prs):
    """한국 주요 기업"""
    slide = setup_slide(prs, "[한국] 주요 연료전지 기업 생태계", 16)

    tbl_shape = slide.shapes.add_table(8, 3, Inches(0.6), Inches(1.5),
                                       Inches(12.1), Inches(5.2))
    table = tbl_shape.table
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(5.0)
    table.columns[2].width = Inches(4.6)

    style_header_row(table, ["기업", "핵심 역할", "주요 실적"])
    data = [
        ["현대자동차", "FCEV (넥쏘, XCIENT), HTWO 시스템",
         "넥쏘 2세대 110kW 스택\n2030년 70만기, 3세대 200kW(2027)"],
        ["두산퓨얼셀", "PAFC M400(440kW) + SOFC",
         "SOFC 양산(군산 50MW, 620°C)\nPAFC 국산화율 98%"],
        ["블룸SK퓨얼셀", "SOFC 국내 생산 (블룸에너지 JV)",
         "구미 공장 준공, 효율 65%\nDC 330kW 설치(국내 최초)"],
        ["SK그룹", "수소 생산-유통-충전 밸류체인",
         "인천 3만톤 액화수소(7,000억원)\n충전소 17개→30개+(2026)"],
        ["효성그룹", "충전시스템 1위 + 탄소섬유 + 액화수소",
         "린데 세계최대 액화수소 1.3만톤\n탄소섬유 1조원(2만4천톤)"],
        ["한화솔루션", "그린수소 (PEM/AEM 수전해)",
         "평창 연 290톤 그린수소\n2.8조원 투자(태양광+수전해)"],
        ["범한퓨얼셀", "PEMFC (잠수함AIP, 선박, 충전소)",
         "장보고III 100% 국산화 납품\n2025년 흑자 전환 전망"],
    ]
    style_data_rows(table, data, size=11)


def slide_17_korea_rd(prs):
    """한국 R&D"""
    slide = setup_slide(prs, "[한국] 연료전지 R&D 핵심 성과", 17)

    tbl_shape = slide.shapes.add_table(7, 4, Inches(0.6), Inches(1.5),
                                       Inches(12.1), Inches(4.0))
    table = tbl_shape.table
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(5.0)
    table.columns[2].width = Inches(2.8)
    table.columns[3].width = Inches(2.3)

    style_header_row(table, ["기관", "연구 내용", "게재지", "시기"])
    data = [
        ["KAIST", "연료전지 촉매 열화 과정 원자 단위 3D 추적 (세계 최초)",
         "Nature Comm.", "2025.8"],
        ["KAIST", "백금-아연 나노입자 촉매 (백금 사용량 1/3 절감)",
         "Chem. Eng. J.", "2025.2"],
        ["KAIST", "이리듐 나노시트 촉매 (상용 대비 성능 13배 향상)",
         "ACS Nano", "2025.12"],
        ["KIST", "친환경 수소 연료전지 성능·효율 증대 핵심 원리 규명",
         "-", "2025"],
        ["KIER", "고효율 연료전지용 1kW급 스택 제작",
         "-", "2024-2025"],
        ["두산퓨얼셀", "중저온(600°C) SOFC 양산 | 61.7% 발전효율 8kW SOFC",
         "KGS 인증", "2025"],
    ]
    style_data_rows(table, data, size=12)

    # 하단 특허 정보
    tf = add_body_textbox(slide, top=Inches(5.8), height=Inches(0.8))
    add_bullet(tf, "특허: 수전해 기술 — 중국 30% > 독일 20% > 일본 18% > 미국 11% > 한국 10% (5위)",
               first=True, size=14, color=MEDIUM_GRAY)
    add_bullet(tf, "한국 PCT 국제특허출원 세계 4위 (5년 연속)", size=14, color=MEDIUM_GRAY)


def slide_18_korea_infra(prs):
    """한국 인프라"""
    slide = setup_slide(prs, "[한국] 수소 인프라 현황", 18)

    # 좌측: 충전소
    add_colored_box(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.7),
                    NAVY, "수소충전소 현황", text_size=18)

    tbl_shape = slide.shapes.add_table(5, 2, Inches(0.6), Inches(2.3),
                                       Inches(5.8), Inches(2.5))
    table = tbl_shape.table
    table.columns[0].width = Inches(3.0)
    table.columns[1].width = Inches(2.8)
    style_header_row(table, ["구분", "현황/목표"])
    data = [
        ["2025.3 운영 중", "407기"],
        ["2027년 목표", "550기"],
        ["2030년 목표", "660기"],
        ["2040년 목표", "1,200기"],
    ]
    style_data_rows(table, data)

    # 우측: 수소 생산
    add_colored_box(slide, Inches(6.8), Inches(1.5), Inches(5.9), Inches(0.7),
                    NAVY, "수소 생산 현황", text_size=18)

    tf = add_body_textbox(slide, left=Inches(7.0), top=Inches(2.4), width=Inches(5.5), height=Inches(2.5))
    add_bullet(tf, "전체 수소 생산의 90%+ = 그레이수소", first=True, size=16, bold=True, color=ACCENT_RED)
    add_bullet(tf, "천연가스 개질: 수소 1톤 → CO2 10톤", level=1, size=14)
    add_bullet(tf, "그린수소 가격: 5,200~6,500원/kg", size=16, bold=True)
    add_bullet(tf, "한화솔루션: 평창 연 290톤 생산", level=1, size=14)
    add_bullet(tf, "수전해 기지: 전북 부안, 강원 평창", level=1, size=14)

    # 하단: 공급 목표
    add_colored_box(slide, Inches(0.6), Inches(5.2), Inches(12.1), Inches(0.7),
                    LIGHT_NAVY, "수소 공급 목표: 2030년 390만톤 → 2040년 526만톤 → 2050년 2,790만톤",
                    text_size=16)

    # 저장/운송
    tf2 = add_body_textbox(slide, top=Inches(6.1), height=Inches(0.7))
    add_bullet(tf2, "저장·운송: 효성-린데 세계최대 액화수소 플랜트 | K-조선 액화수소 운반선 개발 | 탄소섬유 수소탱크",
               first=True, size=14, color=MEDIUM_GRAY)


def slide_19_korea_swot(prs):
    """한국 강점 약점"""
    slide = setup_slide(prs, "[한국] 연료전지 분야 강점과 약점", 19)

    # 강점 (좌측)
    add_colored_box(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.7),
                    GREEN, "강점 (Strengths)", text_size=20)

    tf1 = add_body_textbox(slide, left=Inches(0.6), top=Inches(2.3), width=Inches(5.8), height=Inches(4.2))
    strengths = [
        "발전용 연료전지 세계 1위 (1GW+)",
        "넥쏘 수소차 세계 1위 (42.9%), 4만대+",
        "세계 최초 수소법 + 청정수소 입찰시장",
        "5대 그룹 43.4조원 대규모 투자",
        "KAIST 등 세계적 R&D (Nature 게재)",
        "생산-저장-운송-활용 통합 밸류체인",
        "XCIENT 유럽 165대, 2,000만km 돌파",
        "수소버스 2024년 1,000대+ 신규 보급",
    ]
    for i, s in enumerate(strengths):
        add_bullet(tf1, f"  {s}", first=(i == 0), size=14, color=DARK_GRAY, space_after=Pt(4))

    # 약점 (우측)
    add_colored_box(slide, Inches(6.8), Inches(1.5), Inches(5.9), Inches(0.7),
                    ACCENT_RED, "약점 (Weaknesses)", text_size=20)

    tf2 = add_body_textbox(slide, left=Inches(6.8), top=Inches(2.3), width=Inches(5.9), height=Inches(4.2))
    weaknesses = [
        "전해질막(80%+), 촉매(85%+) 수입 의존",
        "CHPS 2024 낙찰 11.5%, 2025 입찰 취소",
        "수소 생산의 90%+가 그레이수소",
        "충전소 대부분 적자(HyNet 4년 166억 손실)",
        "수소가격 10,239원/kg, 경제성 부족",
        "수전해·수소액화 선도국 대비 5~7년 격차",
        "중국 추격 가속(특허 69%, 부품국산화 70%)",
        "로드맵 달성률: 수소차 23%, 발전 13.5%",
    ]
    for i, w in enumerate(weaknesses):
        add_bullet(tf2, f"  {w}", first=(i == 0), size=14, color=DARK_GRAY, space_after=Pt(4))


def slide_20_roadmap_achievement(prs):
    """로드맵 달성률 분석"""
    slide = setup_slide(prs, "[한국] 수소경제 로드맵 달성률 분석", 20)

    # 상단 설명
    tf = add_body_textbox(slide, top=Inches(1.4), height=Inches(0.6))
    add_bullet(tf, "2019년 로드맵 목표 대비 2025년 현재 실제 달성 현황 — 분야별 상당한 편차 존재",
               first=True, size=16, bold=True, color=NAVY)

    # 달성률 테이블
    tbl_shape = slide.shapes.add_table(7, 5, Inches(0.6), Inches(2.1),
                                       Inches(12.1), Inches(3.8))
    table = tbl_shape.table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(2.2)
    table.columns[3].width = Inches(1.5)
    table.columns[4].width = Inches(4.0)

    style_header_row(table, ["분야", "2022년 목표", "실제 달성", "달성률", "2030년 목표 대비 현황"])
    data = [
        ["수소차 보급", "8.1만 대", "~1.9만 대", "23%", "18만 대 목표 → 현재 3.8만 대"],
        ["수소충전소", "310개소", "~170개소(→407기)", "55%→74%", "660기 목표 → 양호한 진척"],
        ["발전용 연료전지", "1.5GW", "~1.0GW", "67%", "8GW 목표 → 1.08GW (13.5%)"],
        ["건물용 연료전지", "50MW", "~13MW", "26%", "에네팜 49만대 대비 현저히 부족"],
        ["수소 공급량", "47만 톤", "~22만 톤", "47%", "390만 톤 목표 → 그린 전환 시급"],
        ["수소버스", "-", "2,066대", "-", "2024년 1,000대+ (277% 급증)"],
    ]
    style_data_rows(table, data, size=12)

    # 하단: 달성률 평가 박스
    add_colored_box(slide, Inches(0.6), Inches(6.1), Inches(3.8), Inches(0.7),
                    GREEN, "양호: 충전소, 수소버스", text_size=14)
    add_colored_box(slide, Inches(4.6), Inches(6.1), Inches(3.8), Inches(0.7),
                    ACCENT_ORANGE, "보통: 발전용, 수소공급", text_size=14)
    add_colored_box(slide, Inches(8.6), Inches(6.1), Inches(4.1), Inches(0.7),
                    ACCENT_RED, "저조: 수소차(23%), 건물용(26%)", text_size=14)


def slide_21_parts_techgap(prs):
    """부품 국산화 + 기술격차"""
    slide = setup_slide(prs, "[한국] 핵심부품 국산화율 및 기술격차 분석", 21)

    # 좌측: 부품 국산화율 테이블
    add_colored_box(slide, Inches(0.5), Inches(1.4), Inches(6.2), Inches(0.6),
                    NAVY, "PEMFC 핵심부품 국산화 현황", text_size=16)

    tbl1 = slide.shapes.add_table(6, 3, Inches(0.5), Inches(2.1),
                                  Inches(6.2), Inches(3.3))
    t1 = tbl1.table
    t1.columns[0].width = Inches(2.0)
    t1.columns[1].width = Inches(1.5)
    t1.columns[2].width = Inches(2.7)
    style_header_row(t1, ["핵심부품", "국산화율", "해외 의존 / 국내 기업"])
    data1 = [
        ["전해질막 (PEM)", "10~20%", "Chemours(Nafion) / 코오롱"],
        ["촉매 (Pt/C)", "10~15%", "JM, Umicore / 연구단계"],
        ["MEA (막전극접합체)", "20~30%", "Gore, 3M / FCMT, 코오롱"],
        ["GDL (가스확산층)", "50~60%", "SGL, Toray / 제이앤티지"],
        ["분리판 (금속)", "60~70%", "국내 강점 / 케이퓨얼셀"],
    ]
    style_data_rows(t1, data1, size=11)

    # 우측: 기술격차 테이블
    add_colored_box(slide, Inches(7.0), Inches(1.4), Inches(5.7), Inches(0.6),
                    ACCENT_RED, "선도국 대비 기술격차 (NIGT 평가)", text_size=16)

    tbl2 = slide.shapes.add_table(7, 3, Inches(7.0), Inches(2.1),
                                  Inches(5.7), Inches(3.3))
    t2 = tbl2.table
    t2.columns[0].width = Inches(2.2)
    t2.columns[1].width = Inches(1.3)
    t2.columns[2].width = Inches(2.2)
    style_header_row(t2, ["기술 분야", "격차", "선도국"])
    data2 = [
        ["PEMFC 스택", "1~3년", "일본(도요타)"],
        ["SOFC 발전용", "3~5년", "미국(Bloom)"],
        ["수전해 (PEM)", "3~5년", "독일, 미국"],
        ["전해질막", "5~7년", "미국(Chemours)"],
        ["촉매", "5~7년", "영국(JM)"],
        ["수소터빈", "10년+", "미국(GE)"],
    ]
    style_data_rows(t2, data2, size=11)

    # 하단: 중국 추격 경고
    tf = add_body_textbox(slide, top=Inches(5.6), height=Inches(1.2))
    add_bullet(tf, "중국 추격 경고: 연료전지 특허 글로벌 69% 장악 | 부품 국산화 70% (한국 20~30%)",
               first=True, size=15, bold=True, color=ACCENT_RED)
    add_bullet(tf, "SynStack GIII 4.5+ kW/L (현대 2.5세대 ~3.5 kW/L) | 비용 연간 33% 하락 추세",
               size=14, color=DARK_GRAY)


def slide_22_charging_crisis(prs):
    """충전소 수익성 위기"""
    slide = setup_slide(prs, "[한국] 수소충전소 수익성 위기 분석", 22)

    # 좌측: 수익성 현황 테이블
    add_colored_box(slide, Inches(0.5), Inches(1.4), Inches(6.2), Inches(0.6),
                    ACCENT_RED, "충전소 수익성 현황 — 구조적 적자", text_size=16)

    tbl_shape = slide.shapes.add_table(7, 2, Inches(0.5), Inches(2.1),
                                       Inches(6.2), Inches(3.5))
    table = tbl_shape.table
    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(3.0)
    style_header_row(table, ["항목", "현황"])
    data = [
        ["충전소 1기 설치비", "30~50억 원"],
        ["가동률", "20~25% (손익분기 미달)"],
        ["하루 평균 충전", "4대에 불과"],
        ["흑자 충전소 수", "전국 약 7곳"],
        ["수소 판매가 (2025.7)", "10,239원/kg"],
        ["정부 지원 (152곳)", "연 82억원 (충전소당 5,400만원)"],
    ]
    style_data_rows(table, data, size=12)

    # 우측: HyNet 적자 추이
    add_colored_box(slide, Inches(7.0), Inches(1.4), Inches(5.7), Inches(0.6),
                    NAVY, "HyNet(한국수소충전소) 적자 추이", text_size=16)

    tbl2 = slide.shapes.add_table(6, 2, Inches(7.0), Inches(2.1),
                                  Inches(5.7), Inches(3.0))
    t2 = tbl2.table
    t2.columns[0].width = Inches(2.5)
    t2.columns[1].width = Inches(3.2)
    style_header_row(t2, ["연도", "적자액"])
    data2 = [
        ["2019년", "11억 4,000만 원"],
        ["2020년", "22억 5,800만 원"],
        ["2021년", "58억 8,200만 원"],
        ["2022년", "84억 5,000만 원"],
        ["4년 누적", "166억 원 (639% 급증)"],
    ]
    style_data_rows(t2, data2, size=12)

    # 하단: 원인 + 수소가격 추이
    tf = add_body_textbox(slide, top=Inches(5.8), height=Inches(1.2))
    add_bullet(tf, "적자 원인: 수소차 4만대 대비 407기 충전소 → 충전소당 100대, 일 4대 충전",
               first=True, size=14, bold=True, color=DARK_GRAY)
    add_bullet(tf, "수소 가격: 8,000원(2021) → 9,000원(2023, 러-우전쟁) → 10,000원+(2024) — 경유보다 비쌈",
               size=14, color=ACCENT_RED)


def slide_23_tech_innovation(prs):
    """기술 혁신 트렌드"""
    slide = setup_slide(prs, "최근 기술 혁신 및 비용 절감 트렌드", 23)

    # 상단: 촉매 혁신
    tf = add_body_textbox(slide, top=Inches(1.5), height=Inches(2.0))
    add_bullet(tf, "촉매 기술 혁신 — 백금 의존도 탈피", first=True, size=18, bold=True, color=NAVY)
    add_bullet(tf, "1990년대 대비 kW당 백금 사용량 90% 감소 | DOE 목표: 0.10 g/kW 미만", level=1, size=15)
    add_bullet(tf, "PGM-free 촉매: Fe-N-C 0.85 W/cm² | 코발트 기반 4배 내구성 | Fe-Cu 302 mW/cm²", level=1, size=15)
    add_bullet(tf, "워싱턴대 (2026.2): CVD 기법으로 철 촉매 내구성 획기적 개선", level=1, size=15)

    # 하단: 비용 절감 테이블
    tbl_shape = slide.shapes.add_table(5, 3, Inches(0.6), Inches(3.8),
                                       Inches(8.5), Inches(2.8))
    table = tbl_shape.table
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.5)

    style_header_row(table, ["항목", "현재 수준", "목표"])
    data = [
        ["수소 생산 비용", "$4~6/kg", "$1/kg (2031)"],
        ["연료전지 시스템 (수송용)", "$80~100/kW", "$80/kW (2030)"],
        ["전해조 시스템", "~$400/kW", "$250/kW (2026)"],
        ["백금 촉매 비용", "$6.80/kW", "$4.18/kW"],
    ]
    style_data_rows(table, data, size=14)

    # 우측 박스
    add_colored_box(slide, Inches(9.5), Inches(3.8), Inches(3.2), Inches(2.8),
                    NAVY,
                    "AI 데이터센터\n전력 혁명\n\n2030년까지\nDC 38%가\n온사이트 발전\n\n"
                    "SOFC가\n그리드 우회 솔루션",
                    text_size=14, bold=False)


def slide_24_korea_global(prs):
    """한국 글로벌 위상"""
    slide = setup_slide(prs, "[한국] 글로벌 시장에서의 위상", 24)

    tbl_shape = slide.shapes.add_table(8, 3, Inches(0.6), Inches(1.5),
                                       Inches(12.1), Inches(4.2))
    table = tbl_shape.table
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(3.0)
    table.columns[2].width = Inches(5.6)

    style_header_row(table, ["분야", "한국의 위상", "비고"])
    data = [
        ["발전용 연료전지 설치량", "세계 1위", "1,036MW (1GW 돌파)"],
        ["수소전기차 판매", "세계 1위 (현대차)", "글로벌 점유율 42.9%"],
        ["수소차 보유 대수", "세계 2위", "누적 4만대 (중국 다음)"],
        ["수소법 제정", "세계 최초", "2020년 시행"],
        ["청정수소발전 입찰시장", "세계 최초", "2024년 개설"],
        ["수전해 특허", "세계 5위", "전체의 10%"],
        ["수소충전소 수", "세계 3~4위권", "407기 (2025.3)"],
    ]
    style_data_rows(table, data, size=13)

    # 하단 핵심 메시지
    tf = add_body_textbox(slide, top=Inches(6.0), height=Inches(0.8))
    add_bullet(tf, "수소차 + 발전용 연료전지 양 축에서 세계 선두 | 단, 수전해·저장·운송·핵심소재는 선도국 대비 7년+ 격차",
               first=True, size=15, bold=True, color=NAVY)


def slide_25_conclusion(prs):
    """결론 및 전망"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, NAVY)

    # 상단 녹색 라인
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Inches(0.6), SLIDE_WIDTH, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()

    # 제목
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "결론 및 향후 전망"
    set_font(run, size=40, bold=True, color=WHITE)

    # 글로벌 전망
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = "글로벌: 연평균 20~27% 성장, 2030년 수백억 달러 규모"
    set_font(run2, size=20, bold=True, color=GREEN)

    items = [
        ("AI 데이터센터", "SOFC 대규모 배치가 시장 판도 변화의 게임체인저"),
        ("상용차 확대", "수소 트럭·버스 보급 가속화, PEMFC 수요 견인"),
        ("한국의 기회", "2040년 연 43조원 부가가치, 42만개 일자리 창출"),
    ]

    for i, (title, desc) in enumerate(items):
        y = Inches(3.0) + Inches(i * 0.9)
        add_colored_box(slide, Inches(0.8), y, Inches(2.8), Inches(0.7),
                        GREEN, title, text_size=16, bold=True)
        txB = slide.shapes.add_textbox(Inches(4.0), y, Inches(8.5), Inches(0.7))
        tfB = txB.text_frame
        tfB.word_wrap = True
        pB = tfB.paragraphs[0]
        pB.alignment = PP_ALIGN.LEFT
        runB = pB.add_run()
        runB.text = desc
        set_font(runB, size=18, color=WHITE)

    # 한국 핵심 과제
    txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.5))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.LEFT
    run3 = p3.add_run()
    run3.text = "한국 5대 핵심 과제"
    set_font(run3, size=20, bold=True, color=ACCENT_ORANGE)

    challenges = [
        "핵심 소재·부품 국산화",
        "그린수소 전환 가속화",
        "정책 안정성 확보",
        "인프라 확충",
        "기술 격차 해소",
    ]
    box_w = Inches(2.2)
    for i, ch in enumerate(challenges):
        add_colored_box(slide, Inches(0.8 + i * 2.4), Inches(6.3),
                        Inches(2.2), Inches(0.7),
                        LIGHT_NAVY, ch, text_size=13, bold=True)

    add_slide_number(slide, TOTAL_SLIDES)


# ── 메인 ──────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_01_cover(prs)
    slide_02_toc(prs)
    slide_03_global_market(prs)
    slide_04_fc_types(prs)
    slide_05_applications(prs)
    slide_06_global_policy(prs)
    slide_07_us_policy(prs)
    slide_08_us_companies(prs)
    slide_09_us_applications(prs)
    slide_10_china_market(prs)
    slide_11_china_companies(prs)
    slide_12_china_applications(prs)
    slide_13_us_vs_china(prs)
    slide_14_korea_market(prs)
    slide_15_korea_policy(prs)
    slide_16_korea_companies(prs)
    slide_17_korea_rd(prs)
    slide_18_korea_infra(prs)
    slide_19_korea_swot(prs)
    slide_20_roadmap_achievement(prs)
    slide_21_parts_techgap(prs)
    slide_22_charging_crisis(prs)
    slide_23_tech_innovation(prs)
    slide_24_korea_global(prs)
    slide_25_conclusion(prs)

    prs.save(OUTPUT_PATH)
    print(f"PPT 생성 완료: {OUTPUT_PATH}")
    print(f"총 {TOTAL_SLIDES}장 슬라이드")


if __name__ == "__main__":
    main()
